from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT_PPTX = OUT_DIR / "2026-08-12_HCCJP76_論理構成図・物理構成図.pptx"

# PowerPoint標準の日本語資料に馴染む游ゴシックへ統一する。
# run.font.name だけでは日本語（East Asian）の書体指定にならず、LibreOfficeや
# PowerPointが別フォントへフォールバックするため、a:ea も明示的に設定する。
FONT = "Yu Gothic"
NAVY = RGBColor(0x0B, 0x14, 0x26)
AZURE = RGBColor(0x00, 0x78, 0xD4)
CYAN = RGBColor(0x00, 0xB7, 0xC3)
GREEN = RGBColor(0x10, 0x7C, 0x10)
RED = RGBColor(0xD1, 0x34, 0x38)
ORANGE = RGBColor(0xD8, 0x3B, 0x01)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
GRAY = RGBColor(0x60, 0x67, 0x70)
MID_GRAY = RGBColor(0xA1, 0xA7, 0xAE)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
PALE_BLUE = RGBColor(0xE8, 0xF3, 0xFC)
PALE_GREEN = RGBColor(0xE9, 0xF5, 0xEA)
PALE_ORANGE = RGBColor(0xFE, 0xF0, 0xE7)
PALE_PURPLE = RGBColor(0xF3, 0xEA, 0xF7)


def add_text(slide, x, y, w, h, text, size=16, color=BLACK, bold=False,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.06):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    for index, line in enumerate(text.split("\n")):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run_props = run._r.get_or_add_rPr()
        for tag_name in ("a:latin", "a:ea", "a:cs"):
            old = run_props.find(qn(tag_name))
            if old is not None:
                run_props.remove(old)
            font_tag = OxmlElement(tag_name)
            font_tag.set("typeface", FONT)
            run_props.append(font_tag)
    return shape


def add_box(slide, x, y, w, h, title, body, fill, line, title_color=WHITE,
            body_color=BLACK, title_size=17, body_size=13, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    title_h = min(0.55, h * 0.32)
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(title_h)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = line
    header.line.fill.background()
    add_text(slide, x + 0.05, y + 0.01, w - 0.10, title_h - 0.02, title,
             size=title_size, color=title_color, bold=True)
    add_text(slide, x + 0.10, y + title_h + 0.05, w - 0.20, h - title_h - 0.10,
             body, size=body_size, color=body_color)
    return shape


def add_container(slide, x, y, w, h, title, fill, line, title_color=BLACK,
                  title_size=14, dash=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = 12
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    if dash:
        shape.line.dash_style = 4
    add_text(slide, x + 0.12, y + 0.04, w - 0.24, 0.32, title,
             size=title_size, color=title_color, bold=True, align=PP_ALIGN.LEFT)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=AZURE, width=2.3,
              dashed=False, begin=False, end=True):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    if dashed:
        connector.line.dash_style = 4
    line = connector._element.find(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}ln"
    )
    if line is not None:
        if begin:
            line.append(line.makeelement(qn("a:headEnd"), {"type": "triangle"}))
        if end:
            line.append(line.makeelement(qn("a:tailEnd"), {"type": "triangle"}))
    return connector


def add_title(slide, title, subtitle):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.88)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_text(slide, 0.45, 0.07, 8.6, 0.46, title, size=25, color=WHITE,
             bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, 8.85, 0.12, 4.05, 0.34, subtitle, size=11, color=WHITE,
             align=PP_ALIGN.RIGHT)


def build_logical(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    add_title(slide, "論理構成 — 「監視が赤い」から復旧確認まで", "HCCJP 第76回 / 2026-08-14")
    add_text(slide, 0.65, 1.02, 12.0, 0.48,
             "障害は視聴者が選ぶ。AIは対象・症状・原因を知らない。", size=17,
             color=NAVY, bold=True)

    xs = [0.45, 3.02, 5.59, 8.16, 10.73]
    centers = [x + 1.08 for x in xs]
    for left, right in zip(centers[:-1], centers[1:]):
        add_arrow(slide, left + 1.08, 3.03, right - 1.08, 3.03, color=AZURE, width=2.8)

    add_box(slide, xs[0], 2.02, 2.16, 2.05, "① 検知", "Azure Monitor\n可用性テストが失敗\n🔴 赤", PALE_ORANGE, RED)
    add_box(slide, xs[1], 2.02, 2.16, 2.05, "② 判断", "AIエージェント\n証拠を収集し\n原因を切り分ける", PALE_PURPLE, PURPLE)
    add_box(slide, xs[2], 2.02, 2.16, 2.05, "③ 実行", "Azure Arc\naz CLI + Run Command\nRBACで範囲を限定", PALE_BLUE, AZURE)
    add_box(slide, xs[3], 2.02, 2.16, 2.05, "④ 復旧", "オンプレミス\nWindows / Linux\nIIS / nginx", LIGHT, GRAY)
    add_box(slide, xs[4], 2.02, 2.16, 2.05, "⑤ 確認", "同じ可用性テストで\nHTTP 200を再確認\n🟢 緑", PALE_GREEN, GREEN)

    add_text(slide, 0.62, 4.35, 2.0, 0.42, "アラート＋監視証拠", size=11, color=AZURE, bold=True)
    add_text(slide, 3.15, 4.35, 2.0, 0.42, "仮説 → 検証", size=11, color=AZURE, bold=True)
    add_text(slide, 5.70, 4.35, 2.0, 0.42, "調査・修正コマンド", size=11, color=AZURE, bold=True)
    add_text(slide, 8.25, 4.35, 2.0, 0.42, "復旧後も同じ物差し", size=11, color=AZURE, bold=True)

    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(5.12), Inches(12.08), Inches(1.32)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_text(slide, 0.85, 5.25, 11.62, 0.35,
             "人はサーバーへ直接ログインしない", size=19, color=WHITE, bold=True)
    add_text(slide, 0.85, 5.72, 11.62, 0.42,
             "VPNなし　|　踏み台なし　|　インバウンド開放なし　|　Activity Log＋AI実行ログで追跡",
             size=13, color=WHITE)
    add_text(slide, 0.50, 6.83, 12.33, 0.35,
             "ポイント：成功そのものより、証拠を取ってから動くか／失敗時に次の仮説へ進めるかを見る",
             size=12, color=GRAY, align=PP_ALIGN.LEFT)


def build_physical(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    add_title(slide, "Nested Hyper-Vラボの物理構成", "Azure × Cloudflare / 実設定＋当日予定")

    # Boundaries first so every data path remains visible above the pale backgrounds.
    add_container(slide, 0.25, 1.10, 4.35, 5.72, "Microsoft Azure（Japan East）", PALE_BLUE, AZURE, AZURE, 16)
    add_container(slide, 4.82, 1.52, 1.30, 2.22, "公開経路", PALE_ORANGE, ORANGE, ORANGE, 13)
    add_container(slide, 6.32, 1.10, 6.76, 5.72, "オンプレミス検証環境", LIGHT, GRAY, NAVY, 16)
    add_container(slide, 6.62, 1.68, 6.16, 4.85, "L0：物理Hyper-Vホスト（唯一の前提）", WHITE, NAVY, NAVY, 13)
    add_container(slide, 8.45, 1.98, 4.05, 4.25, "L1：nested-lab-01（Nested Hyper-V）", PALE_BLUE, NAVY, NAVY, 12)
    add_container(slide, 8.72, 2.91, 3.53, 2.70, "L2：LabNAT 10.10.0.0/24", WHITE, CYAN, CYAN, 11)

    # Data paths are added before node boxes so they never obscure node text.
    add_arrow(slide, 1.29, 2.96, 1.29, 3.18, color=GREEN, width=2.3, end=False)
    add_arrow(slide, 1.29, 3.18, 5.47, 3.18, color=GREEN, width=2.3, end=False)
    add_arrow(slide, 5.47, 3.18, 5.47, 3.44, color=GREEN, width=2.3)
    add_arrow(slide, 6.00, 3.17, 9.15, 3.43, color=ORANGE, width=2.0)
    add_arrow(slide, 6.00, 3.36, 10.89, 3.43, color=ORANGE, width=2.0)
    add_arrow(slide, 1.42, 5.03, 2.62, 2.96, color=PURPLE, width=2.3)

    # Arc management path runs under the nested boxes, then branches upward to both L2 guests.
    add_arrow(slide, 3.48, 2.96, 3.48, 6.13, color=AZURE, width=2.2, end=False)
    add_arrow(slide, 3.48, 6.13, 11.57, 6.13, color=AZURE, width=2.2, end=False)
    add_arrow(slide, 9.47, 6.13, 9.47, 5.28, color=AZURE, width=2.2, begin=True)
    add_arrow(slide, 11.17, 6.13, 11.17, 5.28, color=AZURE, width=2.2, begin=True)

    # Build-only route stays entirely inside the on-prem boundary.
    add_arrow(slide, 8.20, 5.25, 8.92, 4.82, color=MID_GRAY, width=1.4, dashed=True)

    # Nodes and nested details.
    add_box(slide, 0.50, 1.72, 1.58, 1.24, "Azure Monitor", "App Insights\n可用性テスト / Alert\n5分間隔", PALE_BLUE, AZURE, body_size=9, title_size=12)
    add_box(slide, 2.62, 1.72, 1.58, 1.24, "Azure Arc", "Control Plane\nRun Command\nActivity Log", PALE_BLUE, AZURE, body_size=9, title_size=12)
    add_box(slide, 2.62, 3.52, 1.58, 1.34, "Arcリソース", "arcwin01\narclnx01\nrg-hccjp76-arc", LIGHT, AZURE, body_size=9, title_size=12)
    add_box(slide, 0.50, 5.03, 1.58, 1.30, "AI操作端末", "Claude Code\nAzure CLI（az）\n原因は未通知", PALE_PURPLE, PURPLE, body_size=9, title_size=12)
    add_box(slide, 4.96, 2.02, 1.02, 1.42, "Cloudflare", "当日予定\nNamed Tunnel\n公開URL\n受信FW不要", PALE_ORANGE, ORANGE, body_size=7, title_size=10)
    add_box(slide, 6.87, 4.84, 1.33, 1.20, "制御VM", "Ubuntu + Ansible\n10.20.0.10\n構築時のみ", PALE_GREEN, GREEN, body_size=8, title_size=11)
    add_text(slide, 8.72, 2.40, 3.48, 0.30,
             "Windows Server 2025 / 8 vCPU / 32GB / 10.20.0.20",
             size=9, color=NAVY, bold=True)
    add_text(slide, 8.72, 2.66, 3.48, 0.26,
             "LabNAT Gateway 10.10.0.1",
             size=9, color=NAVY)
    add_box(slide, 8.90, 3.43, 1.25, 1.85, "arcwin01", "Windows Server 2025\nIIS\n10.10.0.51\n2 vCPU / 6GB", PALE_BLUE, AZURE, body_size=8, title_size=11)
    add_box(slide, 10.64, 3.43, 1.25, 1.85, "arclnx01", "Ubuntu 24.04\nnginx\n10.10.0.41\n2 vCPU / 4GB", PALE_GREEN, GREEN, body_size=8, title_size=11)

    # Labels live in reserved whitespace rather than on top of nodes.
    add_text(slide, 1.55, 2.95, 3.20, 0.28, "監視：HTTPS GET / HTTP 200＋固有文字列", size=8, color=GREEN, bold=True)
    add_text(slide, 0.75, 3.80, 1.82, 0.38, "az / HTTPS 443", size=9, color=PURPLE, bold=True)
    add_text(slide, 5.92, 3.54, 2.64, 0.36, "cloudflared：アウトバウンド 443", size=8, color=ORANGE, bold=True)
    add_text(slide, 4.25, 5.72, 4.70, 0.34,
             "Arc agent：アウトバウンド HTTPS 443 / Run Command",
             size=9, color=AZURE, bold=True)
    add_text(slide, 6.55, 4.18, 2.08, 0.44, "構築時のみ\nAnsible / WinRM / SSH", size=8, color=GRAY, bold=True)

    add_text(slide, 0.40, 6.92, 12.55, 0.32,
             "青＝Arc管理経路　緑＝監視HTTP　橙＝Cloudflare Tunnel　灰点線＝ラボ構築経路（ライブ復旧では使用しない）",
             size=10, color=GRAY, align=PP_ALIGN.LEFT)


def build_ai_operation_path(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    add_title(slide, "AIからオンプレミスへ届くまで", "Azure Arcによる操作経路")
    add_text(slide, 0.65, 1.05, 12.0, 0.46,
             "AIはオンプレミスへ直接接続しない。Azureのコントロールプレーンを経由する。",
             size=17, color=NAVY, bold=True)

    # Draw arrows first so card faces cover their endpoints cleanly.
    add_arrow(slide, 2.85, 3.20, 3.55, 3.20, color=PURPLE, width=3.2)
    add_arrow(slide, 6.00, 3.20, 6.70, 3.20, color=AZURE, width=3.2)
    add_arrow(slide, 9.15, 3.20, 9.85, 3.20, color=AZURE, width=3.2)

    add_box(slide, 0.45, 2.02, 2.40, 2.35,
            "① 操作PC", "AIエージェント\nClaude Code\nAzure CLI（az）",
            PALE_PURPLE, PURPLE, title_size=18, body_size=15)
    add_box(slide, 3.55, 2.02, 2.45, 2.35,
            "② Microsoft Azure", "認証・権限確認\n対象リソースの特定\n操作要求を受け付ける",
            PALE_BLUE, AZURE, title_size=17, body_size=14)
    add_box(slide, 6.70, 2.02, 2.45, 2.35,
            "③ Azure Arc", "Control Plane\nRun Command\nActivity Log",
            PALE_BLUE, AZURE, title_size=18, body_size=15)

    add_container(slide, 9.85, 1.70, 3.03, 3.02,
                  "④ オンプレミス", LIGHT, GRAY, NAVY, 16)
    add_box(slide, 10.10, 2.35, 1.17, 1.82,
            "Windows", "Windows Server\nIIS\nArc agent",
            PALE_BLUE, AZURE, title_size=12, body_size=9)
    add_box(slide, 11.46, 2.35, 1.17, 1.82,
            "Linux", "Ubuntu\nnginx\nArc agent",
            PALE_GREEN, GREEN, title_size=12, body_size=9)

    add_text(slide, 2.78, 2.48, 0.86, 0.42, "az CLI", size=10,
             color=PURPLE, bold=True)
    add_text(slide, 5.93, 2.48, 0.86, 0.42, "操作要求", size=10,
             color=AZURE, bold=True)
    add_text(slide, 9.02, 2.48, 0.96, 0.42, "Run Command", size=9,
             color=AZURE, bold=True)

    path_band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.65), Inches(5.08), Inches(12.02), Inches(1.10)
    )
    path_band.fill.solid()
    path_band.fill.fore_color.rgb = NAVY
    path_band.line.fill.background()
    add_text(slide, 0.90, 5.22, 11.52, 0.34,
             "オンプレ側からAzureへ接続を確立", size=18,
             color=WHITE, bold=True)
    add_text(slide, 0.90, 5.66, 11.52, 0.30,
             "Connected Machine agent  →  アウトバウンド HTTPS 443  →  受信ポート開放なし",
             size=13, color=WHITE)

    add_text(slide, 0.58, 6.58, 12.18, 0.40,
             "AIの居場所は操作PC。Azure Arcが、Azure上の操作要求をオンプレミスのWindows / Linuxへ届ける。",
             size=13, color=GRAY, bold=True)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_logical(prs)
    build_physical(prs)
    build_ai_operation_path(prs)
    prs.core_properties.title = "HCCJP第76回 構成図3枚"
    prs.core_properties.subject = "Azure Arc × AIによる無人障害復旧"
    prs.core_properties.author = "胡田 昌彦 / AI共同作成"
    prs.save(OUT_PPTX)
    print(OUT_PPTX)


if __name__ == "__main__":
    build()
