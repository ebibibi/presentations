#!/usr/bin/env python3
"""marp2pptx.py - Marp Markdown → PowerPoint with auto-sizing fonts.

Usage:
    python3 marp2pptx.py slides.md [-o output.pptx]

Features:
    - Parses Marp Markdown (frontmatter, --- separators, class directives)
    - Auto-sizes body text to maximize readability without overflow
    - Clean Azure-themed design
    - Supports: headings, bullets, tables, code blocks, images, tags
"""

import argparse
import math
import re
import sys
import html as html_mod
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════════
# Constants
# ═══════════════════════════════════════════════════════════════

SLIDE_W = 13.33
SLIDE_H = 7.5
MARGIN = 0.55
TITLE_H_MIN = 0.7
TITLE_H_MAX = 1.8
ACCENT_H = 0.05
PAGE_NUM_H = 0.4

# Colors
AZURE = RGBColor(0x00, 0x78, 0xD4)
AZURE_DARK = RGBColor(0x24, 0x3A, 0x5E)
TEXT_COLOR = RGBColor(0x32, 0x31, 0x30)
GRAY = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HDR_BG = "E8F0FE"

TAG_COLORS = {
    "onpre": RGBColor(0x2E, 0x7D, 0x32),
    "cloud": RGBColor(0x15, 0x65, 0xC0),
    "saas": RGBColor(0xC6, 0x28, 0x28),
}

FONT = "Meiryo"
MONO = "Consolas"

FONT_RANGES = {
    "lead": (48, 18, 32),   # (title_max, body_min, body_max)
    "point": (54, 20, 42),
    "": (40, 14, 36),
    "small": (36, 12, 28),
    "x-small": (32, 11, 24),
    "xx-small": (28, 10, 22),
}


# ═══════════════════════════════════════════════════════════════
# Text measurement
# ═══════════════════════════════════════════════════════════════

def display_width(text):
    """Full-width=2, half-width=1."""
    return sum(2 if unicodedata.east_asian_width(c) in ("F", "W") else 1 for c in text)


def clean(text):
    """Strip markdown/HTML for measurement."""
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    t = re.sub(r"\*(.+?)\*", r"\1", t)
    t = re.sub(r"`(.+?)`", r"\1", t)
    t = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", t)
    t = re.sub(r"<br\s*/?>", "\n", t)
    t = re.sub(r"<[^>]+>", "", t)
    return html_mod.unescape(t)


def est_lines(text, font_pt, avail_w_pt):
    """Estimate wrapped line count."""
    total = 0
    for line in clean(text).split("\n"):
        dw = display_width(line.strip())
        cpl = max(1, avail_w_pt / (font_pt * 0.55))
        total += max(1, math.ceil(dw / cpl))
    return total


# ═══════════════════════════════════════════════════════════════
# Markdown Parser
# ═══════════════════════════════════════════════════════════════

def parse_marp(text):
    """Parse Marp markdown into list of slide dicts."""
    fm = re.match(r"^---\n(.+?\n)---\n", text, re.DOTALL)
    if fm:
        text = text[fm.end():]
    raw_slides = re.split(r"\n---\n", text)
    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if raw:
            slides.append(_parse_slide(raw))
    return slides


def _parse_slide(raw):
    slide = {"class": "", "elements": [], "bg_image": None, "bg_pos": "", "bg_pct": 30}

    # Extract class
    m = re.search(r"<!--\s*_class:\s*(\S+)\s*-->", raw)
    if m:
        slide["class"] = m.group(1)
        raw = raw[: m.start()] + raw[m.end():]

    # Extract bg image
    m = re.search(r"!\[bg\s+([^\]]*)\]\(([^)]+)\)", raw)
    if m:
        attrs, slide["bg_image"] = m.group(1), m.group(2)
        pm = re.search(r"(right|left):(\d+)%", attrs)
        if pm:
            slide["bg_pos"] = pm.group(1)
            slide["bg_pct"] = int(pm.group(2))
        raw = raw[: m.start()] + raw[m.end():]

    lines = raw.strip().split("\n")
    i = 0
    elems = slide["elements"]

    while i < len(lines):
        line = lines[i]
        s = line.strip()

        if not s or s.startswith("<!--"):
            i += 1
            continue

        # Heading
        hm = re.match(r"^(#{1,3})\s+(.+)", s)
        if hm:
            elems.append({"type": f"h{len(hm.group(1))}", "text": hm.group(2)})
            i += 1
            continue

        # HTML table
        if s.startswith("<table"):
            tlines = [s]
            i += 1
            while i < len(lines) and "</table>" not in tlines[-1]:
                tlines.append(lines[i])
                i += 1
            elems.append({"type": "html_table", "html": "\n".join(tlines)})
            continue

        # Markdown table
        if "|" in s and s.startswith("|"):
            rows = []
            while i < len(lines) and "|" in lines[i].strip():
                r = lines[i].strip()
                if re.match(r"^\|[\s\-:|]+\|$", r):
                    i += 1
                    continue
                cells = [c.strip() for c in r.split("|")[1:-1]]
                rows.append(cells)
                i += 1
            if rows:
                elems.append({"type": "md_table", "rows": rows})
            continue

        # Arch-box div
        if 'class="arch-box"' in s:
            code_lines = []
            i += 1
            in_code = False
            while i < len(lines) and "</div>" not in lines[i]:
                if lines[i].strip().startswith("```"):
                    in_code = not in_code
                    i += 1
                    continue
                if in_code:
                    code_lines.append(lines[i])
                i += 1
            i += 1  # skip </div>
            elems.append({"type": "code", "text": "\n".join(code_lines), "arch": True})
            continue

        # Code block
        if s.startswith("```"):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            if i < len(lines):
                i += 1
            elems.append({"type": "code", "text": "\n".join(code_lines), "arch": False})
            continue

        # Bullet list
        if re.match(r"^[-*]\s|^\d+\.\s", s):
            items = []
            while i < len(lines):
                bm = re.match(r"^(\s*)([-*]|\d+\.)\s+(.+)", lines[i])
                if bm:
                    level = 1 if len(bm.group(1)) >= 2 else 0
                    items.append({"text": bm.group(3), "level": level})
                    i += 1
                elif lines[i].strip() and items and not re.match(r"^[#|<`-]", lines[i].strip()):
                    items[-1]["text"] += " " + lines[i].strip()
                    i += 1
                else:
                    break
            elems.append({"type": "bullets", "items": items})
            continue

        # Paragraph
        plines = [s]
        i += 1
        while i < len(lines):
            ns = lines[i].strip()
            if not ns or re.match(r"^[#|<`-]|^\d+\.", ns) or ns.startswith("<!--"):
                break
            plines.append(ns)
            i += 1
        elems.append({"type": "para", "text": "\n".join(plines)})

    return slide


def _parse_html_table(html_str):
    """Parse HTML table → list of rows, each row = list of {text, header, colspan}."""
    rows = []
    for tr in re.finditer(r"<tr>(.*?)</tr>", html_str, re.DOTALL):
        cells = []
        for cm in re.finditer(r"<(th|td)([^>]*)>(.*?)</\1>", tr.group(1), re.DOTALL):
            is_hdr = cm.group(1) == "th"
            attrs = cm.group(2)
            txt = clean(cm.group(3)).strip()
            cs_m = re.search(r'colspan="(\d+)"', attrs)
            cs = int(cs_m.group(1)) if cs_m else 1
            cells.append({"text": txt, "header": is_hdr, "colspan": cs})
        rows.append(cells)
    return rows


# ═══════════════════════════════════════════════════════════════
# Auto-sizing
# ═══════════════════════════════════════════════════════════════

def _table_row_heights(rows_data, tfont, total_w_pt, ncols):
    """Estimate height per row based on cell text wrapping.

    Returns list of heights in points, one per row.
    """
    col_w_pt = total_w_pt / max(ncols, 1)
    CELL_PAD = 10  # top+bottom padding in points
    row_heights = []
    for row in rows_data:
        max_lines = 1
        ci = 0
        for cell in row:
            cs = cell.get("colspan", 1) if isinstance(cell, dict) else 1
            txt = cell.get("text", cell) if isinstance(cell, dict) else str(cell)
            cell_w = col_w_pt * cs - 10  # subtract cell margins
            lines = est_lines(txt, tfont, cell_w)
            max_lines = max(max_lines, lines)
            ci += cs
        row_heights.append(max_lines * tfont * 1.5 + CELL_PAD)
    return row_heights


def _content_height(elems, bfont, w_pt):
    """Estimate total content height at body font size (in points).

    Intentionally overestimates (~120%) to prevent overflow.
    TEXT_TO_FIT_SHAPE handles the remaining slack.
    """
    h = 0.0
    ELEMENT_GAP = 8  # points between elements
    for e in elems:
        t = e["type"]
        if t in ("h1", "h2", "h3"):
            lvl = int(t[1])
            sz = bfont * {1: 1.5, 2: 1.2, 3: 1.0}[lvl]
            h += est_lines(e["text"], sz, w_pt) * sz * 1.6 + sz * 0.5
        elif t == "bullets":
            for item in e["items"]:
                h += est_lines(item["text"], bfont, w_pt - 30) * bfont * 1.7
            h += bfont * 0.4
        elif t == "html_table":
            rows_data = _parse_html_table(e["html"])
            ncols = max((sum(c["colspan"] or 1 for c in r) for r in rows_data), default=1)
            tfont = int(bfont * 0.65)
            rh = _table_row_heights(rows_data, tfont, w_pt, ncols)
            h += sum(rh) + 15
        elif t == "md_table":
            rows = e.get("rows", [])
            ncols = max((len(r) for r in rows), default=1)
            tfont = int(bfont * 0.7)
            # Convert to dict format for _table_row_heights
            rows_data = [
                [{"text": c, "colspan": 1} for c in r] for r in rows
            ]
            rh = _table_row_heights(rows_data, tfont, w_pt, ncols)
            h += sum(rh) + 15
        elif t == "code":
            cf = bfont * 0.6
            nlines = e["text"].count("\n") + 1
            h += nlines * cf * 1.4 + 30  # padding inside code box
        elif t == "para":
            h += est_lines(e["text"], bfont, w_pt) * bfont * 1.7
        h += ELEMENT_GAP
    return h


def best_font(elems, w_pt, h_pt, lo=14, hi=36):
    """Find largest body font that fits."""
    best = lo
    for sz in range(hi, lo - 1, -1):
        if _content_height(elems, sz, w_pt) <= h_pt:
            return sz
    return best


def _title_height(text, font_pt, avail_w_inches):
    """Estimate title height in inches based on text and font size."""
    w_pt = avail_w_inches * 72
    lines = est_lines(text, font_pt, w_pt)
    h_pt = lines * font_pt * 1.35 + font_pt * 0.3  # line height + padding
    h_in = h_pt / 72
    return max(TITLE_H_MIN, min(h_in, TITLE_H_MAX))


# ═══════════════════════════════════════════════════════════════
# Inline text formatting
# ═══════════════════════════════════════════════════════════════

_INLINE_RE = re.compile(
    r"(\*\*(.+?)\*\*)"                              # bold
    r"|(\*(.+?)\*)"                                  # italic
    r"|(`([^`]+)`)"                                  # code
    r"|(<b>(.+?)</b>)"                               # html bold
    r'|(<span\s+class="speaker">(.+?)</span>)'       # speaker
    r'|(<span\s+class="tag-(\w+)">(.+?)</span>)'     # tag
    r"|(<span[^>]*>(.+?)</span>)"                     # other span
    r"|(\[([^\]]+)\]\([^)]+\))",                      # link
    re.DOTALL,
)


def _add_runs(para, text, font_pt, color, bold=False):
    """Add inline-formatted runs to a paragraph."""
    text = text.replace("<br>", "\n").replace("<br/>", "\n").replace("<br />", "\n")

    pos = 0
    for m in _INLINE_RE.finditer(text):
        if m.start() > pos:
            _run(para, text[pos : m.start()], font_pt, color, bold)
        if m.group(2) is not None:
            _run(para, m.group(2), font_pt, color, True)
        elif m.group(4) is not None:
            _run(para, m.group(4), font_pt, color, bold, italic=True)
        elif m.group(6) is not None:
            _run(para, m.group(6), int(font_pt * 0.85), RGBColor(0xC7, 0x25, 0x4E), mono=True)
        elif m.group(8) is not None:
            _run(para, m.group(8), font_pt, color, True)
        elif m.group(10) is not None:
            _run(para, m.group(10), int(font_pt * 0.75), GRAY)
        elif m.group(13) is not None:
            tag = m.group(12)
            _run(para, m.group(13), int(font_pt * 0.85), TAG_COLORS.get(tag, color), True)
        elif m.group(15) is not None:
            _run(para, m.group(15), font_pt, color, bold)
        elif m.group(17) is not None:
            _run(para, m.group(17), font_pt, color, bold)
        pos = m.end()
    if pos < len(text):
        _run(para, text[pos:], font_pt, color, bold)


def _run(para, text, font_pt, color, bold=False, italic=False, mono=False):
    """Add a single run."""
    text = html_mod.unescape(text)
    if not text:
        return
    run = para.add_run()
    run.text = text
    font_name = MONO if mono else FONT
    run.font.name = font_name
    run.font.size = Pt(font_pt)
    run.font.color.rgb = color
    # Set East Asian font — without this, CJK chars use PowerPoint's default
    _set_ea_font(run, font_name)
    if bold:
        run.font.bold = True
    if italic:
        run.font.italic = True


# ═══════════════════════════════════════════════════════════════
# PPTX Generation
# ═══════════════════════════════════════════════════════════════

def create_pptx(slides, output, src_dir):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    total = len(slides)

    for idx, sd in enumerate(slides):
        sl = prs.slides.add_slide(blank)
        sc = sd["class"]
        if sc == "lead":
            _build_centered(sl, sd, src_dir, idx + 1, total)
        else:
            _build_standard(sl, sd, src_dir, idx + 1, total)

    prs.save(output)
    return total


def _build_centered(sl, sd, src_dir, pn, total):
    sc = sd["class"]
    _accent_bar(sl)

    cw = SLIDE_W - 2 * MARGIN
    cl = MARGIN
    if sd["bg_image"] and sd["bg_pos"]:
        _add_image(sl, sd, src_dir)
        ratio = sd["bg_pct"] / 100
        if sd["bg_pos"] == "right":
            cw *= 1 - ratio - 0.02
        else:
            cl += cw * ratio + 0.1
            cw *= 1 - ratio - 0.02

    t_max, _, b_max = FONT_RANGES.get(sc, FONT_RANGES[""])
    tf = sl.shapes.add_textbox(
        Inches(cl), Inches(MARGIN + 0.3), Inches(cw), Inches(SLIDE_H - 2 * MARGIN - PAGE_NUM_H)
    ).text_frame
    tf.word_wrap = True
    _set_anchor(tf, "ctr")

    first = True
    for e in sd["elements"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.CENTER

        if e["type"].startswith("h"):
            lvl = int(e["type"][1])
            sz = {1: t_max, 2: int(t_max * 0.7), 3: int(t_max * 0.55)}[lvl]
            _add_runs(p, e["text"], sz, AZURE if lvl <= 2 else TEXT_COLOR, bold=(lvl == 1))
            p.space_after = Pt(8)
        elif e["type"] == "bullets":
            for j, item in enumerate(e["items"]):
                if j > 0:
                    p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                prefix = "    • " if item.get("level", 0) else "• "
                _add_runs(p, prefix + item["text"], int(b_max * 0.7), TEXT_COLOR)
                p.space_after = Pt(4)
        elif e["type"] == "para":
            _add_runs(p, e["text"], int(b_max * 0.7), TEXT_COLOR)
            p.space_after = Pt(6)
        elif e["type"] == "code":
            _add_runs(p, e["text"], int(b_max * 0.5), TEXT_COLOR)

    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _page_num(sl, pn, total)


def _build_standard(sl, sd, src_dir, pn, total):
    sc = sd["class"]
    _accent_bar(sl)

    cw = SLIDE_W - 2 * MARGIN
    cl = MARGIN
    if sd["bg_image"] and sd["bg_pos"]:
        _add_image(sl, sd, src_dir)
        ratio = sd["bg_pct"] / 100
        if sd["bg_pos"] == "right":
            cw *= 1 - ratio - 0.02
        else:
            cl += cw * ratio + 0.1
            cw *= 1 - ratio - 0.02

    t_max, b_min, b_max = FONT_RANGES.get(sc, FONT_RANGES[""])

    # Separate title (first h1) from body
    title_elem = None
    body_elems = []
    for e in sd["elements"]:
        if e["type"] == "h1" and title_elem is None:
            title_elem = e
        else:
            body_elems.append(e)

    cur_top = MARGIN + ACCENT_H + 0.1

    # Title
    if title_elem:
        title_h = _title_height(title_elem["text"], t_max, cw)
        tf = sl.shapes.add_textbox(
            Inches(cl), Inches(cur_top), Inches(cw), Inches(title_h)
        ).text_frame
        tf.word_wrap = True
        _add_runs(tf.paragraphs[0], title_elem["text"], t_max, AZURE, bold=True)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        cur_top += title_h + 0.05
        # Accent underline
        _rect(sl, Inches(cl), Inches(cur_top), Inches(cw), Pt(2), AZURE)
        cur_top += 0.15

    # Find optimal body font
    body_top = cur_top
    avail_h = SLIDE_H - body_top - MARGIN - PAGE_NUM_H
    w_pt = cw * 72
    h_pt = avail_h * 72
    bfont = best_font(body_elems, w_pt, h_pt, b_min, b_max)
    # Slight boost for text-only slides (TEXT_TO_FIT_SHAPE handles slack).
    # Skip boost when tables exist — tables auto-expand and can't shrink.
    has_table = any(e["type"] in ("md_table", "html_table") for e in body_elems)
    if not has_table:
        bfont = min(bfont + 2, b_max)

    # Group elements: consecutive text vs table/code
    groups = _group_elements(body_elems)

    for gtype, gelems in groups:
        remaining = SLIDE_H - cur_top - MARGIN - PAGE_NUM_H
        if remaining < 0.2:
            break

        if gtype == "text":
            cur_top = _place_text_group(sl, gelems, bfont, cl, cur_top, cw, remaining)
        elif gtype == "md_table":
            cur_top = _place_md_table(sl, gelems[0], bfont, cl, cur_top, cw, remaining)
        elif gtype == "html_table":
            cur_top = _place_html_table(sl, gelems[0], bfont, cl, cur_top, cw, remaining)
        elif gtype == "code":
            cur_top = _place_code(sl, gelems[0], bfont, cl, cur_top, cw, remaining)

    _page_num(sl, pn, total)


def _group_elements(elems):
    """Group consecutive text elements; tables/code break the group."""
    groups = []
    cur = []
    for e in elems:
        if e["type"] in ("md_table", "html_table", "code"):
            if cur:
                groups.append(("text", cur))
                cur = []
            groups.append((e["type"], [e]))
        else:
            cur.append(e)
    if cur:
        groups.append(("text", cur))
    return groups


def _place_text_group(sl, elems, bfont, left, top, width, max_h):
    """Place a group of text elements in one text frame."""
    est_h = _content_height(elems, bfont, width * 72) / 72
    h = min(max(0.5, est_h + 0.3), max_h)

    tf = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(h)).text_frame
    tf.word_wrap = True

    first = True
    for e in elems:
        if e["type"].startswith("h"):
            lvl = int(e["type"][1])
            sz = int(bfont * {1: 1.4, 2: 1.15, 3: 1.0}[lvl])
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            _add_runs(p, e["text"], sz, AZURE if lvl <= 2 else TEXT_COLOR, bold=(lvl <= 2))
            p.space_after = Pt(int(sz * 0.3))
        elif e["type"] == "bullets":
            for item in e["items"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                lvl = item.get("level", 0)
                prefix = "    • " if lvl else "• "
                _add_runs(p, prefix + item["text"], bfont, TEXT_COLOR)
                p.space_after = Pt(int(bfont * 0.25))
        elif e["type"] == "para":
            for line in e["text"].split("\n"):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                _add_runs(p, line, bfont, TEXT_COLOR)
                p.space_after = Pt(int(bfont * 0.2))

    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return top + h + 0.1


def _place_md_table(sl, elem, bfont, left, top, width, max_h):
    rows = elem["rows"]
    if not rows:
        return top
    ncols = max(len(r) for r in rows)
    nrows = len(rows)
    tfont = int(bfont * 0.7)
    rows_data = [[{"text": c, "colspan": 1} for c in r] for r in rows]
    rh = _table_row_heights(rows_data, tfont, width * 72, ncols)
    h = min(max(0.6, sum(rh) / 72 + 0.2), max_h)

    tbl = sl.shapes.add_table(nrows, ncols, Inches(left), Inches(top), Inches(width), Inches(h)).table

    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            if ci >= ncols:
                break
            cell = tbl.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            is_hdr = ri == 0
            _add_runs(p, clean(cell_text), tfont, AZURE_DARK if is_hdr else TEXT_COLOR, bold=is_hdr)
            if is_hdr:
                _cell_bg(cell, TABLE_HDR_BG)

    return top + h + 0.1


def _place_html_table(sl, elem, bfont, left, top, width, max_h):
    rows_data = _parse_html_table(elem["html"])
    if not rows_data:
        return top
    ncols = max(sum(c["colspan"] or 1 for c in r) for r in rows_data)
    nrows = len(rows_data)
    tfont = int(bfont * 0.65)
    rh = _table_row_heights(rows_data, tfont, width * 72, ncols)
    h = min(max(0.6, sum(rh) / 72 + 0.2), max_h)

    tbl = sl.shapes.add_table(nrows, ncols, Inches(left), Inches(top), Inches(width), Inches(h)).table

    for ri, row in enumerate(rows_data):
        ci = 0
        for cell_info in row:
            if ci >= ncols:
                break
            cs = cell_info["colspan"]
            if cs > 1 and ci + cs - 1 < ncols:
                tbl.cell(ri, ci).merge(tbl.cell(ri, ci + cs - 1))
            cell = tbl.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            is_hdr = cell_info["header"]
            _add_runs(p, cell_info["text"], tfont, AZURE_DARK if is_hdr else TEXT_COLOR, bold=is_hdr)
            if is_hdr:
                _cell_bg(cell, TABLE_HDR_BG)
            ci += max(cs, 1)

    return top + h + 0.1


def _place_code(sl, elem, bfont, left, top, width, max_h):
    code = elem["text"]
    is_arch = elem.get("arch", False)
    nlines = code.count("\n") + 1

    if is_arch:
        # ASCII art: calculate exact font size to fit width, no auto-fit
        max_dw = max(display_width(line) for line in code.split("\n"))
        avail_w_pt = width * 72 - 20  # subtract padding
        # Monospace: each half-width char ≈ 0.6 × font_size
        cfont = max(8, min(int(avail_w_pt / (max_dw * 0.6)), int(bfont * 0.7)))
        h = min(max(0.6, nlines * cfont * 1.35 / 72 + 0.4), max_h)

        bg = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = CODE_BG
        bg.line.fill.background()

        tf = bg.text_frame
        tf.word_wrap = False  # Never wrap ASCII art
        run = tf.paragraphs[0].add_run()
        run.text = code
        run.font.name = "MS Gothic"
        run.font.size = Pt(cfont)
        run.font.color.rgb = TEXT_COLOR
        _set_ea_font(run, "MS Gothic")
        _set_monospace_props(run)
        # No TEXT_TO_FIT_SHAPE — exact sizing preserves monospace alignment
    else:
        cfont = int(bfont * 0.6)
        h = min(max(0.6, nlines * cfont * 1.4 / 72 + 0.5), max_h)

        bg = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = CODE_BG
        bg.line.fill.background()

        tf = bg.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = code
        run.font.name = MONO
        run.font.size = Pt(cfont)
        run.font.color.rgb = TEXT_COLOR
        _set_ea_font(run, MONO)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    return top + h + 0.1


# ═══════════════════════════════════════════════════════════════
# Shape helpers
# ═══════════════════════════════════════════════════════════════

def _set_ea_font(run, font_name):
    """Set East Asian font on a run (python-pptx only sets Latin)."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", font_name)


def _set_monospace_props(run):
    """Disable kerning and set fixed spacing for ASCII art alignment."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("kern", "0")
    rPr.set("spc", "0")


def _accent_bar(sl):
    s = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W), Inches(ACCENT_H))
    s.fill.solid()
    s.fill.fore_color.rgb = AZURE
    s.line.fill.background()


def _rect(sl, left, top, width, height, color):
    s = sl.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def _page_num(sl, num, total):
    tf = sl.shapes.add_textbox(
        Inches(SLIDE_W - 1.5), Inches(SLIDE_H - 0.4), Inches(1.2), Inches(0.3)
    ).text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY
    run.font.name = FONT


def _add_image(sl, sd, src_dir):
    img_path = Path(src_dir) / sd["bg_image"]
    if not img_path.exists():
        return
    pct = sd["bg_pct"] / 100
    pos = sd["bg_pos"]
    max_w = Inches(SLIDE_W * pct)
    max_h = Inches(SLIDE_H - 1.0)
    if pos == "right":
        img_left = Inches(SLIDE_W * (1 - pct))
    else:
        img_left = Inches(0)
    try:
        # Add with width only to preserve aspect ratio
        pic = sl.shapes.add_picture(str(img_path), img_left, Inches(0.5), width=max_w)
        # If resulting height exceeds max, re-add with height constraint instead
        if pic.height > max_h:
            pic.width = int(pic.width * (max_h / pic.height))
            pic.height = max_h
        # Center vertically in available area
        pic.top = Inches(0.5) + (max_h - pic.height) // 2
    except Exception:
        pass


def _set_anchor(tf, anchor):
    body_pr = tf._txBody.find(qn("a:bodyPr"))
    body_pr.set("anchor", anchor)


def _cell_bg(cell, hex_color):
    tc_pr = cell._tc.get_or_add_tcPr()
    fill = tc_pr.makeelement(qn("a:solidFill"), {})
    fill.append(fill.makeelement(qn("a:srgbClr"), {"val": hex_color}))
    tc_pr.append(fill)


# ═══════════════════════════════════════════════════════════════
# Self-check: overlap detection
# ═══════════════════════════════════════════════════════════════

def _check_overlaps(pptx_path):
    """Re-read generated PPTX and check for overlapping shapes."""
    from pptx import Presentation as PrsRead
    prs = PrsRead(pptx_path)
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        shapes = []
        for sh in slide.shapes:
            # Skip tiny decorative elements (accent bars, page numbers)
            w_in = sh.width / 914400
            h_in = sh.height / 914400
            if w_in < 0.3 or h_in < 0.15:
                continue
            shapes.append({
                "name": sh.shape_type,
                "top": sh.top / 914400,
                "left": sh.left / 914400,
                "bottom": (sh.top + sh.height) / 914400,
                "right": (sh.left + sh.width) / 914400,
                "h": h_in,
            })
        # Check pairwise overlaps (only vertical — same column content)
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                a, b = shapes[i], shapes[j]
                # Horizontal overlap check
                h_overlap = a["left"] < b["right"] and b["left"] < a["right"]
                if not h_overlap:
                    continue
                # Vertical overlap check (with 0.05" tolerance)
                v_overlap = a["top"] < b["bottom"] - 0.05 and b["top"] < a["bottom"] - 0.05
                if v_overlap:
                    overlap_amt = min(a["bottom"], b["bottom"]) - max(a["top"], b["top"])
                    if overlap_amt > 0.1:  # >0.1" overlap is meaningful
                        issues.append(
                            f"  スライド{si}: 要素が{overlap_amt:.2f}\"重なり "
                            f"(top={a['top']:.1f}~{a['bottom']:.1f} vs "
                            f"top={b['top']:.1f}~{b['bottom']:.1f})"
                        )
    return issues


# ═══════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(description="Marp MD → PPTX with auto-sizing")
    parser.add_argument("input", help="Input .md file")
    parser.add_argument("-o", "--output", help="Output .pptx (default: same basename)")
    parser.add_argument("--no-check", action="store_true", help="Skip overlap check")
    args = parser.parse_args()

    inp = Path(args.input)
    if not inp.exists():
        print(f"❌ {inp} が見つかりません")
        sys.exit(1)

    out = args.output or str(inp.with_suffix(".pptx"))
    slides = parse_marp(inp.read_text(encoding="utf-8"))
    n = create_pptx(slides, out, inp.parent)
    print(f"✅ {out} ({n}スライド)")

    if not args.no_check:
        issues = _check_overlaps(out)
        if issues:
            print(f"\n⚠️  重なり検出 ({len(issues)}件):")
            for issue in issues:
                print(issue)
        else:
            print("✅ 重なりなし")


if __name__ == "__main__":
    main()
