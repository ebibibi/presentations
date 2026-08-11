#!/usr/bin/env python3
"""マネージドIDのリソースURLを解説するYouTube用PowerPointを生成。"""

import sys
from pathlib import Path

sys.path.insert(0, "/home/ebi/.claude/skills/youtube-slide/scripts")

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from slide_helpers import (
    BG,
    BLUE,
    CARD_BORDER,
    CARD_FILL,
    EMU,
    GREEN,
    LABEL,
    ORANGE,
    PURPLE,
    RED,
    TEAL,
    TEXT_MAIN,
    TEXT_MUTED,
    add_slide,
    box,
    header,
    new_deck,
    text,
    title_slide,
)


OUT_DIR = Path("/home/ebi/wt-1536588859213619210/20260811_ManagedIdentity_ResourceURL")
OUT_FILE = OUT_DIR / "マネージドIDのリソースURLを理解する.pptx"


def line(slide, x1, y1, x2, y2, color=TEXT_MUTED, width=2.0, arrow=False):
    shape = slide.shapes.add_connector(
        1, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2))
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if arrow:
        try:
            shape.line.end_arrowhead = True
        except Exception:
            pass
    return shape


def pill(slide, x, y, w, label, color, size=17):
    box(slide, x, y, w, 0.48 * EMU, fill=color, border=color, radius=0.25)
    text(
        slide,
        x,
        y,
        w,
        0.48 * EMU,
        [{"text": label, "size": size, "bold": True, "color": BG}],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def card(slide, x, y, w, h, title, body, accent=BLUE, title_size=22, body_size=17):
    box(slide, x, y, w, h, fill=CARD_FILL, border=CARD_BORDER, bw=1.2)
    box(slide, x, y, 0.08 * EMU, h, fill=accent, border=accent, radius=0.01)
    text(
        slide,
        x + 0.25 * EMU,
        y + 0.18 * EMU,
        w - 0.45 * EMU,
        0.5 * EMU,
        [{"text": title, "size": title_size, "bold": True, "color": accent}],
    )
    text(
        slide,
        x + 0.25 * EMU,
        y + 0.78 * EMU,
        w - 0.45 * EMU,
        h - 0.92 * EMU,
        [{"text": body, "size": body_size, "color": TEXT_MAIN}],
    )


def big_statement(prs, num, label, top, main, sub, color=ORANGE):
    s = add_slide(prs)
    text(s, 0.48 * EMU, 0.38 * EMU, 0.8 * EMU, 0.4 * EMU,
         [{"text": num, "size": 18, "bold": True, "color": color}])
    text(s, 1.15 * EMU, 0.4 * EMU, 6 * EMU, 0.4 * EMU,
         [{"text": label, "size": 13, "color": LABEL}])
    text(s, 0.6 * EMU, 1.55 * EMU, 12.1 * EMU, 0.55 * EMU,
         [{"text": top, "size": 23, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)
    text(s, 0.5 * EMU, 2.3 * EMU, 12.3 * EMU, 1.45 * EMU,
         [{"text": main, "size": 66, "bold": True, "color": color}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.0 * EMU, 4.35 * EMU, 11.3 * EMU, 1.2 * EMU,
         [{"text": sub, "size": 25, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)
    return s


prs = new_deck()

title_slide(
    prs,
    title="その『リソースURL』、何のURL？",
    subtitle="マネージドIDのトークン宛先を理解する",
    tagline="Azure AI Search の例で腹落ち",
    source_url="https://learn.microsoft.com/ja-jp/connectors/custom-connectors/azure-active-directory-authentication",
    footer="2026年8月11日 ｜ ebisuda.net",
)

s = add_slide(prs)
header(s, "01", "WHY CONFUSING", "似たURLが3つ出てくる")
card(s, 0.55 * EMU, 2.05 * EMU, 3.9 * EMU, 3.6 * EMU,
     "通信先 endpoint", "HTTPリクエストを送る住所\n\nhttps://myservice\n.search.windows.net", BLUE)
card(s, 4.72 * EMU, 2.05 * EMU, 3.9 * EMU, 3.6 * EMU,
     "トークンの宛先", "アクセストークンを\n受け取るAPIの識別子\n\nhttps://search.azure.com", ORANGE)
card(s, 8.89 * EMU, 2.05 * EMU, 3.9 * EMU, 3.6 * EMU,
     "Azure resource ID", "Azure上の個別リソース\n\n/subscriptions/.../\nproviders/Microsoft.Search/...", PURPLE, body_size=15)
text(s, 0.6 * EMU, 6.15 * EMU, 12.1 * EMU, 0.65 * EMU,
     [{"text": "名前が似ている。でも、役割はまったく違う。", "size": 25, "bold": True, "color": TEXT_MAIN}],
     align=PP_ALIGN.CENTER)

big_statement(
    prs,
    "02",
    "THE ANSWER",
    "リソースURLは",
    "APIの住所ではない",
    "アクセストークンの『宛先』を指定する識別子",
    ORANGE,
)

s = add_slide(prs)
header(s, "03", "TOKEN FLOW", "Entra IDに『誰向けか』を伝える")
card(s, 0.55 * EMU, 2.35 * EMU, 2.6 * EMU, 2.15 * EMU,
     "① 呼び出し元", "カスタムコネクタ\nアプリ／マネージドID", BLUE, body_size=16)
card(s, 3.7 * EMU, 2.35 * EMU, 2.6 * EMU, 2.15 * EMU,
     "② Entra ID", "宛先を見て\nトークンを発行", PURPLE, body_size=18)
card(s, 6.85 * EMU, 2.35 * EMU, 2.6 * EMU, 2.15 * EMU,
     "③ Access Token", "aud = 対象API\nAPI用の通行証", ORANGE, title_size=18, body_size=18)
card(s, 10.0 * EMU, 2.35 * EMU, 2.75 * EMU, 2.15 * EMU,
     "④ API", "aud と権限を確認\n一致すれば処理", GREEN, body_size=18)
for x in (3.18, 6.33, 9.48):
    text(s, x * EMU, 3.0 * EMU, 0.5 * EMU, 0.55 * EMU,
         [{"text": "→", "size": 30, "bold": True, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)
text(s, 0.75 * EMU, 5.35 * EMU, 11.85 * EMU, 0.85 * EMU,
     [{"text": "リソースURLは ①→② で指定し、発行トークンの aud に反映される", "size": 25, "bold": True, "color": ORANGE}],
     align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "04", "AUDIENCE", "aud = この通行証を受け取ってよい相手")
box(s, 0.75 * EMU, 2.0 * EMU, 5.2 * EMU, 3.6 * EMU, fill=CARD_FILL, border=ORANGE, bw=1.6)
text(s, 1.05 * EMU, 2.28 * EMU, 4.6 * EMU, 2.9 * EMU,
     [
         {"text": "ACCESS TOKEN", "size": 18, "bold": True, "color": ORANGE},
         {"text": '{', "size": 23, "color": TEXT_MUTED, "space_before": 10},
         {"text": '  "aud": "<対象API>",', "size": 19, "bold": True, "color": TEAL},
         {"text": '}', "size": 23, "color": TEXT_MUTED},
     ])
text(s, 6.5 * EMU, 2.05 * EMU, 5.7 * EMU, 0.7 * EMU,
     [{"text": "受信したAPIが確認", "size": 27, "bold": True, "color": GREEN}])
card(s, 6.5 * EMU, 2.95 * EMU, 5.7 * EMU, 1.1 * EMU,
     "aud が自分向け", "→ 次に権限を確認", GREEN, body_size=16)
card(s, 6.5 * EMU, 4.35 * EMU, 5.7 * EMU, 1.1 * EMU,
     "aud が別API向け", "→ この入口で拒否", RED, body_size=16)
text(s, 0.7 * EMU, 6.2 * EMU, 12 * EMU, 0.55 * EMU,
     [{"text": "audの形式はトークン版やAPIで異なる。大切なのは『受信者』という役割。", "size": 22, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "05", "THE LEARN EXAMPLE", "なぜ management.core.windows.net なのか？")
text(s, 0.7 * EMU, 1.85 * EMU, 12 * EMU, 0.8 * EMU,
     [{"text": "あのLearnは Azure Resource Manager API のチュートリアル", "size": 28, "bold": True, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)
card(s, 0.75 * EMU, 3.0 * EMU, 5.5 * EMU, 2.45 * EMU,
     "サンプルが呼ぶAPI", "Azure Resource Manager\nList all subscriptions", BLUE, body_size=21)
card(s, 7.05 * EMU, 3.0 * EMU, 5.5 * EMU, 2.45 * EMU,
     "Learn指定のリソースURL", "https://management\n.core.windows.net/", ORANGE, body_size=22)
text(s, 0.7 * EMU, 6.0 * EMU, 12 * EMU, 0.75 * EMU,
     [{"text": "整理：この値は『Azure全般』ではなく、サンプルが対象にしたARM向け", "size": 23, "bold": True, "color": ORANGE}], align=PP_ALIGN.CENTER)

big_statement(
    prs,
    "06",
    "ENDPOINT ≠ AUDIENCE",
    "通信先とトークンの宛先は",
    "違っていていい",
    "役割が違う。HTTPリクエストの通信先とは限らない。",
    TEAL,
)

s = add_slide(prs)
header(s, "07", "ARM EXAMPLE", "ARMでは2つのURLが並ぶ")
card(s, 0.55 * EMU, 2.0 * EMU, 3.9 * EMU, 3.7 * EMU,
     "現在のARM REST endpoint", "https://management\n.azure.com/\n\nHTTPリクエストの通信先", BLUE, title_size=18, body_size=18)
card(s, 4.72 * EMU, 2.0 * EMU, 3.9 * EMU, 3.7 * EMU,
     "このLearnのOAuth resource", "https://management\n.core.windows.net/\n\nトークン要求で指定", ORANGE, title_size=17, body_size=18)
card(s, 8.89 * EMU, 2.0 * EMU, 3.9 * EMU, 3.7 * EMU,
     "現在のMI ARM例", "resource=\nhttps://management\n.azure.com/\n\n別フローでは別表記", TEAL, title_size=18, body_size=17)
text(s, 0.65 * EMU, 6.03 * EMU, 12.0 * EMU, 0.72 * EMU,
     [{"text": "補足：core.windows.net は classic管理APIでは実通信先にも使われた。推測で置換しない。", "size": 19, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "08", "AZURE AI SEARCH", "検索するなら audience は Search API")
card(s, 0.65 * EMU, 2.05 * EMU, 3.85 * EMU, 3.75 * EMU,
     "接続先 endpoint", "https://myservice\n.search.windows.net\n\nサービスごとに固有", BLUE, body_size=19)
card(s, 4.75 * EMU, 2.05 * EMU, 3.85 * EMU, 3.75 * EMU,
     "resource / audience", "https://search.azure.com\n\nSearch API全体を表す", ORANGE, title_size=17, body_size=19)
card(s, 8.85 * EMU, 2.05 * EMU, 3.85 * EMU, 3.75 * EMU,
     "v2 scope", "https://search.azure.com/\n.default\n\nCLI・MSAL等の要求形式", TEAL, body_size=18)
text(s, 0.65 * EMU, 6.22 * EMU, 12.0 * EMU, 0.52 * EMU,
     [{"text": "3つは同じ文字列にするものではない", "size": 25, "bold": True, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "09", "TWO PLANES", "『何をするか』で対象APIが変わる")
card(s, 0.75 * EMU, 2.0 * EMU, 5.55 * EMU, 3.85 * EMU,
     "Management plane", "Searchサービスを作る・設定する\n\n呼ぶAPI：Azure Resource Manager\nトークン：ARM向け", PURPLE, body_size=19)
card(s, 7.0 * EMU, 2.0 * EMU, 5.55 * EMU, 3.85 * EMU,
     "Data plane", "インデックス内の文書を検索する\n\n呼ぶAPI：Azure AI Search\nトークン：Search向け", GREEN, body_size=19)
text(s, 0.8 * EMU, 6.23 * EMU, 11.75 * EMU, 0.52 * EMU,
     [{"text": "同じAzure AI Searchでも、操作の入口が違えばリソースURLも違う", "size": 24, "bold": True, "color": ORANGE}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "10", "RESOURCE OR SCOPE", "resource と scope は『要求時』の書式")
card(s, 0.7 * EMU, 2.0 * EMU, 3.85 * EMU, 3.65 * EMU,
     "Entra v1系（ADAL）：resource", "resource / リソースURL\n\n対象APIの識別子を指定", BLUE, title_size=16, body_size=17)
card(s, 4.75 * EMU, 2.0 * EMU, 3.85 * EMU, 3.65 * EMU,
     "Searchの識別子", "https://search.azure.com\n\n末尾記号は入力欄の\n公式例に正確に合わせる", ORANGE, title_size=18, body_size=17)
card(s, 8.8 * EMU, 2.0 * EMU, 3.85 * EMU, 3.65 * EMU,
     "Entra v2系（MSAL）：scope", "scope=\nhttps://search.azure.com/\n.default", TEAL, title_size=16, body_size=18)
text(s, 0.65 * EMU, 6.1 * EMU, 12 * EMU, 0.75 * EMU,
     [{"text": "resource / scope は要求側。aud は発行後の受信者。完全な同義語ではない。", "size": 22, "bold": True, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "11", "DO NOT MIX", "3つの『リソース』を分離する")
rows = [
    ("endpoint", "通信先", "https://myservice.search.windows.net", BLUE),
    ("audience / resource", "トークンの宛先", "https://search.azure.com", ORANGE),
    ("Azure resource ID", "RBACを付ける対象", "/subscriptions/.../providers/Microsoft.Search/...", PURPLE),
]
for i, (name, role, value, color) in enumerate(rows):
    y = (1.95 + i * 1.48) * EMU
    box(s, 0.7 * EMU, y, 11.95 * EMU, 1.08 * EMU, fill=CARD_FILL, border=CARD_BORDER)
    pill(s, 0.95 * EMU, y + 0.3 * EMU, 2.45 * EMU, name, color, size=15)
    text(s, 3.65 * EMU, y + 0.22 * EMU, 2.65 * EMU, 0.65 * EMU,
         [{"text": role, "size": 19, "bold": True, "color": TEXT_MAIN}], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.15 * EMU, y + 0.18 * EMU, 6.0 * EMU, 0.72 * EMU,
         [{"text": value, "size": 16, "color": color}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7 * EMU, 6.5 * EMU, 11.95 * EMU, 0.35 * EMU,
     [{"text": "URLっぽいかどうかではなく、どの欄・どの処理で使う値かを見る", "size": 21, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "12", "HOW TO FIND IT", "正しい値は『対象操作』から探す")
steps = [
    ("1", "何をする？", "作成・設定 / 検索・データ操作", BLUE),
    ("2", "どのAPI？", "ARM / Search / Storage…", PURPLE),
    ("3", "トークン要求", "resource / scope の公式値", ORANGE),
    ("4", "文字列をそのまま", "末尾 / と .default まで一致", GREEN),
]
for i, (num, title_txt, body, color) in enumerate(steps):
    x = (0.55 + i * 3.15) * EMU
    box(s, x, 2.05 * EMU, 2.75 * EMU, 3.5 * EMU, fill=CARD_FILL, border=CARD_BORDER)
    text(s, x, 2.25 * EMU, 2.75 * EMU, 0.65 * EMU,
         [{"text": num, "size": 36, "bold": True, "color": color}], align=PP_ALIGN.CENTER)
    text(s, x + 0.15 * EMU, 3.08 * EMU, 2.45 * EMU, 0.55 * EMU,
         [{"text": title_txt, "size": 21, "bold": True, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)
    text(s, x + 0.2 * EMU, 3.95 * EMU, 2.35 * EMU, 1.15 * EMU,
         [{"text": body, "size": 16, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)
text(s, 0.7 * EMU, 6.05 * EMU, 12 * EMU, 0.72 * EMU,
     [{"text": "RBACの範囲は別物 → Azure resource ID。トークン要求と混ぜない。", "size": 22, "bold": True, "color": ORANGE}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "13", "TROUBLESHOOT", "失敗した場所で切り分ける")
card(s, 0.65 * EMU, 2.0 * EMU, 3.85 * EMU, 3.75 * EMU,
     "400 / 401", "resource / audienceを\n公式値と正確に照合\n末尾スラッシュも確認", BLUE, body_size=18)
card(s, 4.75 * EMU, 2.0 * EMU, 3.85 * EMU, 3.75 * EMU,
     "403 Search", "必要なRBACロールを確認\n割り当て反映は\n最大10分かかる場合あり", ORANGE, body_size=18)
card(s, 8.85 * EMU, 2.0 * EMU, 3.85 * EMU, 3.75 * EMU,
     "番号だけで断定しない", "レスポンス本文を見る\nサービス固有の\n診断手順も確認", RED, title_size=18, body_size=18)
text(s, 0.7 * EMU, 6.2 * EMU, 12 * EMU, 0.52 * EMU,
     [{"text": "エラー番号だけで断定せず、レスポンス本文とサービス固有の診断も見る。", "size": 20, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "14", "MINIMUM EXAMPLE", "Azure AI Searchならこうつながる")
box(s, 0.7 * EMU, 1.85 * EMU, 11.95 * EMU, 1.35 * EMU, fill=CARD_FILL, border=TEAL, bw=1.4)
text(s, 1.0 * EMU, 2.12 * EMU, 11.35 * EMU, 0.8 * EMU,
     [{"text": "az account get-access-token --scope https://search.azure.com/.default", "size": 19, "bold": True, "color": TEAL}], anchor=MSO_ANCHOR.MIDDLE)
text(s, 6.1 * EMU, 3.35 * EMU, 1.0 * EMU, 0.65 * EMU,
     [{"text": "↓", "size": 32, "bold": True, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)
box(s, 0.7 * EMU, 4.05 * EMU, 11.95 * EMU, 1.75 * EMU, fill=CARD_FILL, border=BLUE, bw=1.4)
text(s, 1.0 * EMU, 4.28 * EMU, 11.35 * EMU, 1.25 * EMU,
     [
         {"text": "GET https://myservice.search.windows.net/indexes('myindex')/docs?...", "size": 18, "bold": True, "color": BLUE},
         {"text": "Authorization: Bearer <取得したアクセストークン>", "size": 18, "color": TEXT_MAIN, "space_before": 10},
     ])
text(s, 0.75 * EMU, 6.22 * EMU, 11.85 * EMU, 0.55 * EMU,
     [{"text": "scope と endpoint が違う文字列でも、それが正しい組み合わせ", "size": 23, "bold": True, "color": ORANGE}], align=PP_ALIGN.CENTER)

big_statement(
    prs,
    "15",
    "TAKEAWAY",
    "迷ったら問うのは",
    "誰向けの通行証？",
    "リソースURLは推測しない。対象APIの公式値をそのまま使う。",
    ORANGE,
)

s = add_slide(prs)
header(s, "16", "OFFICIAL SOURCES", "今回の正本")
sources = [
    ("カスタムコネクタのLearn", "learn.microsoft.com/connectors/custom-connectors/azure-active-directory-authentication"),
    ("アクセストークンのaud", "learn.microsoft.com/entra/identity-platform/access-tokens"),
    ("Access token claims", "learn.microsoft.com/entra/identity-platform/access-token-claims-reference"),
    ("Azure REST APIの基礎", "learn.microsoft.com/rest/api/gettingstarted/"),
    ("Searchキーレス接続", "learn.microsoft.com/azure/search/search-get-started-rbac"),
    ("Search REST API", "learn.microsoft.com/rest/api/searchservice/documents/search-get"),
    ("Search audience定数", "github.com/Azure/azure-sdk-for-net/.../SearchAudience.cs"),
]
for i, (name, url) in enumerate(sources):
    y = (1.65 + i * 0.7) * EMU
    text(s, 0.75 * EMU, y, 3.2 * EMU, 0.48 * EMU,
         [{"text": name, "size": 15, "bold": True, "color": TEXT_MAIN}])
    text(s, 3.95 * EMU, y, 8.6 * EMU, 0.48 * EMU,
         [{"text": url, "size": 12, "color": TEAL}])
text(s, 0.75 * EMU, 6.72 * EMU, 11.85 * EMU, 0.35 * EMU,
     [{"text": "値や画面は更新されるため、実装時は最新の公式ページを再確認", "size": 18, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "17", "OFFICIAL SOURCES", "補助ソース：値の正確性と切り分け")
sources2 = [
    ("Managed Identityトークン", "learn.microsoft.com/entra/identity/managed-identities-azure-resources/how-to-use-vm-token"),
    ("ADALとMSALの違い", "learn.microsoft.com/entra/msal/dotnet/how-to/differences-adal-msal-net"),
    ("Audienceの正確一致", "learn.microsoft.com/azure/logic-apps/authenticate-with-managed-identity"),
    ("SearchのManaged Identity", "learn.microsoft.com/azure/search/search-how-to-managed-identities"),
    ("SearchのRBAC診断", "learn.microsoft.com/azure/search/keyless-connections"),
]
for i, (name, url) in enumerate(sources2):
    y = (1.8 + i * 0.92) * EMU
    text(s, 0.75 * EMU, y, 3.3 * EMU, 0.55 * EMU,
         [{"text": name, "size": 17, "bold": True, "color": TEXT_MAIN}])
    text(s, 4.05 * EMU, y, 8.45 * EMU, 0.55 * EMU,
         [{"text": url, "size": 13, "color": TEAL}])
text(s, 0.75 * EMU, 6.72 * EMU, 11.85 * EMU, 0.35 * EMU,
     [{"text": "ARMとSearchでは、操作・認証フロー・クラウドに合う値を選ぶ", "size": 18, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(OUT_FILE)
print(OUT_FILE)
