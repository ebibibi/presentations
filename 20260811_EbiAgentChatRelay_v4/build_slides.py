#!/usr/bin/env python3
"""Generate the Ebi Agent Chat Relay v4 YouTube explainer deck."""

import sys
from pathlib import Path

sys.path.insert(0, "/home/ebi/.claude/skills/youtube-slide/scripts")

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from slide_helpers import (
    BG,
    BLUE,
    CARD_BORDER,
    CARD_FILL,
    EMU,
    GREEN,
    LABEL,
    ORANGE,
    PURPLE,
    RED,
    TEAL,
    TEXT_MAIN,
    TEXT_MUTED,
    add_slide,
    box,
    header,
    new_deck,
    text,
    title_slide,
)


OUT_DIR = Path("/home/ebi/wt-1534221229802782934-presentations-v4/20260811_EbiAgentChatRelay_v4")
OUT_FILE = OUT_DIR / "Ebi Agent Chat Relay v4 - DiscordとTeamsから複数AIを使う.pptx"


def card(slide, x, y, w, h, title, body, accent=BLUE, title_size=22, body_size=18):
    box(slide, x, y, w, h, fill=CARD_FILL, border=CARD_BORDER, bw=1.2)
    box(slide, x, y, 0.08 * EMU, h, fill=accent, border=accent, radius=0.01)
    text(slide, x + 0.24 * EMU, y + 0.16 * EMU, w - 0.42 * EMU, 0.55 * EMU,
         [{"text": title, "size": title_size, "bold": True, "color": accent}])
    text(slide, x + 0.24 * EMU, y + 0.78 * EMU, w - 0.42 * EMU, h - 0.92 * EMU,
         [{"text": body, "size": body_size, "color": TEXT_MAIN}])


def big_statement(prs, num, label, top, main, sub, color=ORANGE, main_size=62):
    slide = add_slide(prs)
    text(slide, 0.48 * EMU, 0.38 * EMU, 0.8 * EMU, 0.4 * EMU,
         [{"text": num, "size": 18, "bold": True, "color": color}])
    text(slide, 1.15 * EMU, 0.4 * EMU, 6 * EMU, 0.4 * EMU,
         [{"text": label, "size": 13, "color": LABEL}])
    text(slide, 0.6 * EMU, 1.4 * EMU, 12.1 * EMU, 0.55 * EMU,
         [{"text": top, "size": 24, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)
    text(slide, 0.5 * EMU, 2.15 * EMU, 12.3 * EMU, 1.55 * EMU,
         [{"text": main, "size": main_size, "bold": True, "color": color}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, 1.0 * EMU, 4.3 * EMU, 11.3 * EMU, 1.2 * EMU,
         [{"text": sub, "size": 25, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)
    return slide


prs = new_deck()

title_slide(
    prs,
    title="AIエージェントをTeamsへ",
    subtitle="Ebi Agent Chat Relay v4.0.0",
    tagline="Discord＋Teams × Claude Code CLI・Codex CLI・Local OpenAI互換・AG-UI",
    source_url="https://github.com/ebibibi/ebi-agent-chat-relay/releases/tag/v4.0.0",
    footer="2026年8月11日 ｜ ebisuda.net",
)

big_statement(
    prs, "01", "V4.0.0 STRUCTURE", "2つのFrontendと", "4つのBackend",
    "会話する場所と、仕事をするAgentを別々に選べる基盤です。", ORANGE,
)

s = add_slide(prs)
header(s, "02", "TWO INDEPENDENT AXES", "FrontendとBackendを分けた")
card(s, 0.7 * EMU, 2.0 * EMU, 5.55 * EMU, 3.75 * EMU,
     "Frontend｜人が話す場所", "Discord\nMicrosoft Teams\n\n人とAgentをつなぐ会話の入口", BLUE, body_size=22)
card(s, 7.05 * EMU, 2.0 * EMU, 5.55 * EMU, 3.75 * EMU,
     "Backend｜実際に働くAgent", "Claude Code CLI\nOpenAI Codex CLI\nLocal OpenAI互換 /v1/responses\nAG-UI HTTP/SSE", TEAL, body_size=19)
text(s, 0.8 * EMU, 6.25 * EMU, 11.7 * EMU, 0.55 * EMU,
     [{"text": "2つのFrontendから、4つのBackendを選べる", "size": 23, "bold": True, "color": TEXT_MAIN}],
     align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "03", "WHY TEAMS IS HARD", "推奨構成はTeams app packageだけではない")
items = [
    ("1", "Entra app", "application", PURPLE),
    ("2", "Azure Bot", "Bot resource", BLUE),
    ("3", "Storage Queue", "transport", ORANGE),
    ("4", "Public Receiver", "検証＋enqueue", TEAL),
    ("5", "Private Host", "outbound pull", GREEN),
]
for i, (n, title_, body, color) in enumerate(items):
    x = (0.45 + i * 2.58) * EMU
    card(s, x, 2.2 * EMU, 2.25 * EMU, 3.5 * EMU, f"{n}  {title_}", body, color, title_size=18, body_size=17)
text(s, 0.8 * EMU, 6.2 * EMU, 11.7 * EMU, 0.55 * EMU,
     [{"text": "推奨経路は、公開受信とprivate実行をQueueで分離", "size": 25, "bold": True, "color": ORANGE}],
     align=PP_ALIGN.CENTER)

big_statement(
    prs, "04", "THE DESIGN RULE", "公開入口は必要。でも", "Agent Hostを公開しない",
    "repository access・agent credentials・Agent実行能力を、公開受信口から分離する。", RED,
)

s = add_slide(prs)
header(s, "05", "OUTBOUND-ONLY PRIVATE HOST", "Public側は検証とenqueueだけ")
flow = [
    ("Teams", "ユーザー", BLUE),
    ("Bot Framework", "Activity配送", PURPLE),
    ("Public Receiver", "検証＋enqueue", ORANGE),
    ("Storage Queue", "Activity待機", TEAL),
    ("ActivityPuller", "外向きpull", GREEN),
    ("Selected Backend", "Agent実行", BLUE),
]
for i, (title_, body, color) in enumerate(flow):
    x = (0.28 + i * 2.16) * EMU
    card(s, x, 2.35 * EMU, 1.82 * EMU, 2.75 * EMU, title_, body, color, title_size=15, body_size=14)
    if i < len(flow) - 1:
        text(s, (2.08 + i * 2.16) * EMU, 3.2 * EMU, 0.35 * EMU, 0.55 * EMU,
             [{"text": "→", "size": 25, "bold": True, "color": TEXT_MUTED}], align=PP_ALIGN.CENTER)
text(s, 0.6 * EMU, 5.7 * EMU, 12.1 * EMU, 0.95 * EMU,
     [{"text": "Public Receiver：bot client secretなし／Agent起動不可", "size": 21, "bold": True, "color": ORANGE},
      {"text": "　Private Host：Teams listenerを公開しない", "size": 21, "bold": True, "color": GREEN}],
     align=PP_ALIGN.CENTER)
text(s, 0.6 * EMU, 6.5 * EMU, 12.1 * EMU, 0.35 * EMU,
     [{"text": "復路：selected backend → Bot Connector → Teams", "size": 17, "color": TEXT_MUTED}],
     align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "06", "2 × 4 COMBINATIONS", "同じRelayで8通り")
labels = ["Claude Code CLI", "Codex CLI", "Local OpenAI互換", "AG-UI"]
colors = [ORANGE, TEAL, GREEN, PURPLE]
for row, frontend in enumerate(["Discord", "Microsoft Teams"]):
    y = (2.05 + row * 2.15) * EMU
    card(s, 0.55 * EMU, y, 2.45 * EMU, 1.55 * EMU, frontend, "Frontend", BLUE, title_size=20, body_size=15)
    for col, (backend, color) in enumerate(zip(labels, colors)):
        x = (3.3 + col * 2.35) * EMU
        body = "/v1/responses\n✓" if backend == "Local OpenAI互換" else "✓"
        card(s, x, y, 2.05 * EMU, 1.55 * EMU, backend, body, color, title_size=14, body_size=16)
text(s, 0.65 * EMU, 6.3 * EMU, 12.0 * EMU, 0.45 * EMU,
     [{"text": "AG-UI＝HTTP/SSEで接続するBackend", "size": 21, "color": TEXT_MUTED}],
     align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "07", "SETUP MAP", "Teams導入ガイドは8セクション")
steps = [
    "1  Entra app", "2  Azure Bot", "3  Storage Queue", "4  Public Receiver",
    "5  Private session host", "6  Teams app package", "7  Upload & consent", "8  3段階で検証",
]
for i, label in enumerate(steps):
    row, col = divmod(i, 4)
    x = (0.58 + col * 3.15) * EMU
    y = (2.0 + row * 2.0) * EMU
    card(s, x, y, 2.75 * EMU, 1.45 * EMU, label, "", [PURPLE, BLUE, TEAL, ORANGE][col], title_size=17)
text(s, 0.65 * EMU, 6.15 * EMU, 12.0 * EMU, 0.55 * EMU,
     [{"text": "詳細手順：docs/teams-setup.md", "size": 26, "bold": True, "color": GREEN}], align=PP_ALIGN.CENTER)

big_statement(
    prs, "08", "REAL DATA BOUNDARY", "Entra appを顧客tenantに置いても", "Tenant登録 ≠ Tenant内処理",
    "Bot Framework・Receiver・Queue・Private Host・Backendまでがデータ経路。", PURPLE, main_size=52,
)

s = add_slide(prs)
header(s, "09", "AG-UI BACKEND", "HTTP/SSE Agentも同じstreamへ")
left = "Run lifecycle\nText streaming\nReasoning\nTool call / result"
right = "URL credential拒否\nRedirect拒否\nSSE frame上限\nTokenを子CLIへ渡さない"
card(s, 0.7 * EMU, 2.0 * EMU, 5.55 * EMU, 3.75 * EMU, "共通eventへ変換", left, TEAL, body_size=22)
card(s, 7.05 * EMU, 2.0 * EMU, 5.55 * EMU, 3.75 * EMU, "境界を明示", right, ORANGE, body_size=22)
text(s, 0.65 * EMU, 6.18 * EMU, 12.0 * EMU, 0.55 * EMU,
     [{"text": "AG-UI eventをRelayの共通streamへ変換", "size": 25, "bold": True, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "10", "HONEST LIMITS", "v4でも同じではない部分")
card(s, 0.65 * EMU, 2.0 * EMU, 3.85 * EMU, 3.8 * EMU,
     "Teams commands", "/backend等は通常queue経路でcommand dispatchしない\n→ configured/global Backend", RED, title_size=19, body_size=18)
card(s, 4.75 * EMU, 2.0 * EMU, 3.85 * EMU, 3.8 * EMU,
     "Teams files", "通常private queue経路は\nfile-consent invokeをbridgeしない", ORANGE, title_size=19, body_size=18)
card(s, 8.85 * EMU, 2.0 * EMU, 3.85 * EMU, 3.8 * EMU,
     "AG-UI advanced", "Durable HITL resume\nstate/activity\nprotobuf/client tools\nは対応機能として未提示", PURPLE, title_size=19, body_size=15)
text(s, 0.65 * EMU, 6.25 * EMU, 12.0 * EMU, 0.45 * EMU,
     [{"text": "未対応を、対応済みのように見せない", "size": 25, "bold": True, "color": GREEN}], align=PP_ALIGN.CENTER)

s = add_slide(prs)
header(s, "11", "PROVEN ON REAL COMPONENTS", "Contractだけでなく実物で往復")
card(s, 0.65 * EMU, 2.0 * EMU, 3.85 * EMU, 3.75 * EMU, "2,536", "local tests\nruff成功／pyright 0 errors", GREEN, title_size=42, body_size=18)
card(s, 4.75 * EMU, 2.0 * EMU, 3.85 * EMU, 3.75 * EMU, "CI", "Python 3.12／3.13\nCodeQL／merge-after", BLUE, title_size=42, body_size=20)
card(s, 8.85 * EMU, 2.0 * EMU, 3.85 * EMU, 3.75 * EMU, "E2E", "Teams → Azure relay\n→ Real Codex → Teams", ORANGE, title_size=42, body_size=20)
text(s, 0.65 * EMU, 6.2 * EMU, 12.0 * EMU, 0.5 * EMU,
     [{"text": "本番検証でもDiscordとTeamsが同時稼働", "size": 25, "bold": True, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)

big_statement(
    prs, "12", "NEXT STEP", "選ぶのはBotの見た目ではなく", "入口とAgentを、分けて選ぶ",
    "Release NotesとTeams Setup Guideから、自分のdeployment境界を決める。", GREEN,
)

s = add_slide(prs)
header(s, "13", "TRY AND FOLLOW", "v4を試す・続きを見る")
card(s, 0.7 * EMU, 1.9 * EMU, 5.55 * EMU, 3.95 * EMU,
     "GitHub Release", "github.com/ebibibi/\nebi-agent-chat-relay/\nreleases/tag/v4.0.0\n\nRelease Notes\nTeams Setup Guide", TEAL, title_size=25, body_size=17)
card(s, 7.05 * EMU, 1.9 * EMU, 5.55 * EMU, 3.95 * EMU,
     "YouTube／note", "高評価・チャンネル登録\n通知オン\n\nnoteで\n設計と導入手順を解説", ORANGE, title_size=25, body_size=19)
text(s, 0.65 * EMU, 6.18 * EMU, 12.0 * EMU, 0.55 * EMU,
     [{"text": "Ebi Agent Chat Relay v4.0.0 — MIT License", "size": 25, "bold": True, "color": TEXT_MAIN}], align=PP_ALIGN.CENTER)

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(OUT_FILE)
print(OUT_FILE)
