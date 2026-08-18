#!/usr/bin/env python3
"""Active Directory入門 2026年版 Part 1 のYouTube用PowerPointを生成。"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, "/home/ebi/.claude/skills/youtube-slide/scripts")

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from slide_helpers import (
    BG,
    BLUE,
    CARD_BORDER,
    CARD_FILL,
    EMU,
    GREEN,
    LABEL,
    ORANGE,
    PURPLE,
    RED,
    TEAL,
    TEXT_MAIN,
    TEXT_MUTED,
    add_slide,
    box,
    header,
    new_deck,
    text,
    title_slide,
)


OUT_DIR = Path(
    "/home/ebi/wt-1539252221055864932/20260818_ActiveDirectory_2026_Part1"
)
OUT_FILE = OUT_DIR / "Active Directory入門 2026年版 Part1.pptx"


def line(slide, x1, y1, x2, y2, color=TEXT_MUTED, width=2.0):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Emu(int(x1)),
        Emu(int(y1)),
        Emu(int(x2)),
        Emu(int(y2)),
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def header(slide, num, label, title, title_color=TEXT_MAIN):
    """2桁番号がLibreOffice変換でも欠けない幅の共通ヘッダー。"""
    text(
        slide,
        0.45 * EMU,
        0.35 * EMU,
        0.85 * EMU,
        0.5 * EMU,
        [{"text": num, "size": 18, "bold": True, "color": ORANGE}],
    )
    text(
        slide,
        1.3 * EMU,
        0.38 * EMU,
        6 * EMU,
        0.45 * EMU,
        [{"text": label, "size": 13, "color": LABEL}],
    )
    text(
        slide,
        0.45 * EMU,
        0.95 * EMU,
        12.4 * EMU,
        0.9 * EMU,
        [{"text": title, "size": 30, "bold": True, "color": title_color}],
    )


def arrow(slide, x, y, direction="→", color=TEXT_MUTED, size=30):
    return text(
        slide,
        x,
        y,
        0.55 * EMU,
        0.55 * EMU,
        [{"text": direction, "size": size, "bold": True, "color": color}],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def pill(slide, x, y, w, label, color, size=16):
    box(slide, x, y, w, 0.48 * EMU, fill=color, border=color, radius=0.25)
    text(
        slide,
        x,
        y,
        w,
        0.48 * EMU,
        [{"text": label, "size": size, "bold": True, "color": BG}],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def card(
    slide,
    x,
    y,
    w,
    h,
    title_value,
    body,
    accent=BLUE,
    title_size=21,
    body_size=18,
):
    box(slide, x, y, w, h, fill=CARD_FILL, border=CARD_BORDER, bw=1.2)
    box(slide, x, y, 0.08 * EMU, h, fill=accent, border=accent, radius=0.01)
    text(
        slide,
        x + 0.25 * EMU,
        y + 0.16 * EMU,
        w - 0.45 * EMU,
        0.54 * EMU,
        [{"text": title_value, "size": title_size, "bold": True, "color": accent}],
    )
    text(
        slide,
        x + 0.25 * EMU,
        y + 0.77 * EMU,
        w - 0.45 * EMU,
        h - 0.9 * EMU,
        [{"text": body, "size": body_size, "color": TEXT_MAIN}],
        anchor=MSO_ANCHOR.TOP,
    )


def big_statement(prs, num, label, top, main, sub, color=ORANGE):
    slide = add_slide(prs)
    text(
        slide,
        0.48 * EMU,
        0.38 * EMU,
        0.8 * EMU,
        0.4 * EMU,
        [{"text": num, "size": 18, "bold": True, "color": color}],
    )
    text(
        slide,
        1.15 * EMU,
        0.4 * EMU,
        6 * EMU,
        0.4 * EMU,
        [{"text": label, "size": 13, "color": LABEL}],
    )
    text(
        slide,
        0.6 * EMU,
        1.45 * EMU,
        12.1 * EMU,
        0.6 * EMU,
        [{"text": top, "size": 24, "color": TEXT_MUTED}],
        align=PP_ALIGN.CENTER,
    )
    text(
        slide,
        0.5 * EMU,
        2.25 * EMU,
        12.3 * EMU,
        1.55 * EMU,
        [{"text": main, "size": 62, "bold": True, "color": color}],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    text(
        slide,
        1.0 * EMU,
        4.35 * EMU,
        11.3 * EMU,
        1.25 * EMU,
        [{"text": sub, "size": 24, "color": TEXT_MAIN}],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return slide


def bottom_note(slide, value, color=TEXT_MUTED, size=19):
    text(
        slide,
        0.7 * EMU,
        6.3 * EMU,
        11.95 * EMU,
        0.55 * EMU,
        [{"text": value, "size": size, "bold": True, "color": color}],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


prs = new_deck()

title_slide(
    prs,
    title="Active Directory 入門",
    subtitle="2026年版 Part 1｜まず把握すべき概念と用語",
    tagline="用語を、地図に置く",
    source_url="https://www.youtube.com/watch?v=lZ8Ps6U_kvY",
    footer="2026年8月18日 ｜ 胡田昌彦 ｜ ebisuda.net",
)

s = add_slide(prs)
header(s, "02", "REMAKE POLICY", "内容の強みは残す。見づらさは残さない")
card(s, 0.65 * EMU, 2.0 * EMU, 3.75 * EMU, 3.7 * EMU, "音", "BGMなし\n声と説明に集中", GREEN, body_size=22)
card(s, 4.78 * EMU, 2.0 * EMU, 3.75 * EMU, 3.7 * EMU, "文字", "大きく、短く\n1枚1メッセージ", BLUE, body_size=22)
card(s, 8.9 * EMU, 2.0 * EMU, 3.75 * EMU, 3.7 * EMU, "色", "純黒×赤文字を避ける\n色だけで区別しない", ORANGE, body_size=20)
bottom_note(s, "濃紺の背景＋白文字＋明るい補助色で統一", TEXT_MAIN, 23)

big_statement(
    prs,
    "03",
    "TODAY'S GOAL",
    "丸暗記ではなく",
    "用語を、地図に置く",
    "初めて聞く言葉も『論理・物理・機能』の視点で眺めれば迷いにくい",
    TEAL,
)

s = add_slide(prs)
header(s, "04", "FIRST DISTINCTION", "Active Directory と Microsoft Entra ID は別物")
card(s, 0.7 * EMU, 1.95 * EMU, 5.65 * EMU, 3.95 * EMU, "Active Directory Domain Services", "Windows Serverで運用\n\nドメイン参加／GPO\nKerberos・NTLM／LDAP\n社内サーバー・PC管理", BLUE, body_size=19)
card(s, 6.98 * EMU, 1.95 * EMU, 5.65 * EMU, 3.95 * EMU, "Microsoft Entra ID", "クラウドのIDサービス\n\nSaaS／条件付きアクセス\nMFA・パスワードレス\nMicrosoft Entra join", PURPLE, body_size=19)
text(s, 5.95 * EMU, 3.18 * EMU, 1.42 * EMU, 0.95 * EMU, [{"text": "併用\nできる", "size": 18, "bold": True, "color": ORANGE}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bottom_note(s, "この動画の『AD』は Active Directory Domain Services（AD DS）", ORANGE, 21)

s = add_slide(prs)
header(s, "05", "LEARNING METAPHOR", "AD DSを3つの役割で覚える")
card(s, 0.65 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "共有名簿", "誰がいる？\nどのPC？\nどのグループ？", TEAL, body_size=22)
card(s, 4.78 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "本人確認", "誰として\nサインインした？", BLUE, body_size=22)
card(s, 8.9 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "ルール配布", "どの設定を\n誰のPCへ？", GREEN, body_size=22)
bottom_note(s, "学習用の比喩：保存・認証・アクセス制御・ポリシー管理を3役に整理", TEXT_MAIN, 19)

s = add_slide(prs)
header(s, "06", "THREE LENSES", "用語を3つの視点で見る")
card(s, 0.65 * EMU, 2.0 * EMU, 3.75 * EMU, 3.8 * EMU, "① 論理構造", "フォレスト\nドメイン\nOU／コンテナ", PURPLE, body_size=22)
card(s, 4.78 * EMU, 2.0 * EMU, 3.75 * EMU, 3.8 * EMU, "② 物理構造", "DC\nサイト／サブネット\nサイトリンク", BLUE, body_size=22)
card(s, 8.9 * EMU, 2.0 * EMU, 3.75 * EMU, 3.8 * EMU, "③ 提供機能", "DNS／LDAP／認証\nGPO／複製\n信頼／復旧", GREEN, body_size=20)
bottom_note(s, "重なり例｜DC＝AD DSを動かすサーバー＋認証応答　複製＝DC間構成＋同期機能", ORANGE, 17)

s = add_slide(prs)
header(s, "07", "DIRECTORY DATA", "管理対象を『オブジェクト』として保存する")
items = [
    ("ユーザー", "名前／所属／\n認証に必要な情報", TEAL),
    ("コンピューター", "端末名／OS／所属", BLUE),
    ("グループ", "メンバー／権限付与", GREEN),
    ("サービスアカウント", "サービス用ID／\nサービス名の登録", PURPLE),
]
for index, (name, body, color) in enumerate(items):
    x = (0.55 + index * 3.16) * EMU
    card(s, x, 2.05 * EMU, 2.8 * EMU, 3.35 * EMU, name, body, color, body_size=18)
bottom_note(s, "オブジェクト = 管理対象 ／ 属性 = その対象が持つ情報項目", TEXT_MAIN, 22)

s = add_slide(prs)
header(s, "08", "DISTINGUISHED NAME", "ディレクトリ内の住所は DN で表す")
box(s, 0.8 * EMU, 2.0 * EMU, 11.7 * EMU, 1.15 * EMU, fill=CARD_FILL, border=TEAL, bw=1.5)
text(s, 1.0 * EMU, 2.25 * EMU, 11.3 * EMU, 0.65 * EMU, [{"text": "CN=佐藤, OU=営業部, DC=corp, DC=example, DC=com", "size": 25, "bold": True, "color": TEAL}], align=PP_ALIGN.CENTER)
labels = [("CN", "オブジェクト名", 1.0, BLUE), ("OU", "組織単位", 4.15, GREEN), ("DC", "DNS名の各要素", 7.3, PURPLE)]
for code, meaning, x, color in labels:
    card(s, x * EMU, 3.75 * EMU, 2.7 * EMU, 1.65 * EMU, code, meaning, color, title_size=22, body_size=17)
bottom_note(s, "画面のフォルダー表示ではなく、階層を持つディレクトリデータ", TEXT_MUTED, 19)

s = add_slide(prs)
header(s, "09", "LOGICAL MODEL", "フォレストが最上位、ドメインがデータ区画")
box(s, 0.65 * EMU, 1.9 * EMU, 12.0 * EMU, 4.2 * EMU, fill=None, border=PURPLE, bw=2.0)
pill(s, 1.0 * EMU, 2.15 * EMU, 2.1 * EMU, "FOREST", PURPLE)
card(s, 1.0 * EMU, 3.0 * EMU, 5.15 * EMU, 2.35 * EMU, "corp.example.com", "フォレストルートドメイン\n自分のディレクトリ区画を持つ", BLUE, body_size=19)
card(s, 7.15 * EMU, 3.0 * EMU, 4.45 * EMU, 2.35 * EMU, "child.corp.example.com", "必要なら子ドメインを追加\n同じフォレストの一員", TEAL, body_size=18)
arrow(s, 6.35 * EMU, 3.95 * EMU, "→", ORANGE, 34)
bottom_note(s, "同じフォレストはスキーマ・構成・グローバルカタログを共有", TEXT_MAIN, 20)

big_statement(
    prs,
    "10",
    "DESIGN PRINCIPLE",
    "フォレストやドメインは",
    "要件が先、構造は後",
    "設計の目安｜『分けられる』と『分けるべき』は別。まず単純な構成から",
    ORANGE,
)

s = add_slide(prs)
header(s, "11", "DOMAIN ≠ DNS", "ADドメインはDNS名を使う。でもDNSそのものではない")
card(s, 0.7 * EMU, 2.05 * EMU, 5.55 * EMU, 3.65 * EMU, "AD DSのドメイン", "ユーザーやPCを持つ\nディレクトリのデータ区画\n\n例：corp.example.com", BLUE, body_size=19)
card(s, 7.05 * EMU, 2.05 * EMU, 5.55 * EMU, 3.65 * EMU, "DNSのドメイン／ゾーン", "名前からIPやサービスを探す\n名前解決の範囲\n\n例：corp.example.com", TEAL, body_size=19)
text(s, 5.85 * EMU, 3.15 * EMU, 1.6 * EMU, 1.0 * EMU, [{"text": "同じ名前\n別の役割", "size": 18, "bold": True, "color": ORANGE}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bottom_note(s, "AD DSはDNSを使ってドメインコントローラーとサービスを見つける", TEXT_MAIN, 21)

s = add_slide(prs)
header(s, "12", "DNS DEPENDENCY", "DNSが壊れると、ADは『場所が分からない』")
card(s, 0.75 * EMU, 2.2 * EMU, 3.1 * EMU, 2.5 * EMU, "クライアント", "corpのDCは\nどこ？", BLUE, body_size=20)
arrow(s, 4.0 * EMU, 3.15 * EMU, "→", TEXT_MUTED, 32)
card(s, 4.55 * EMU, 2.2 * EMU, 4.15 * EMU, 2.5 * EMU, "DC Locator", "DNSでSRVレコードを検索\nサイト等の条件でDCを選択", PURPLE, body_size=17)
arrow(s, 8.85 * EMU, 3.15 * EMU, "→", TEXT_MUTED, 32)
card(s, 9.4 * EMU, 2.2 * EMU, 3.15 * EMU, 2.5 * EMU, "DC", "LDAP／Kerberosへ\n接続", GREEN, body_size=19)
bottom_note(s, "切り分けの目安｜DNS参照先でADドメインのSRVレコードを解決できるか", ORANGE, 19)

s = add_slide(prs)
header(s, "13", "SRV RECORD", "SRVレコード = サービスの場所を知らせるDNSレコード")
box(s, 0.7 * EMU, 1.95 * EMU, 11.95 * EMU, 1.15 * EMU, fill=CARD_FILL, border=TEAL, bw=1.5)
text(s, 0.9 * EMU, 2.25 * EMU, 11.55 * EMU, 0.55 * EMU, [{"text": "_ldap._tcp.dc._msdcs.corp.example.com", "size": 25, "bold": True, "color": TEAL}], align=PP_ALIGN.CENTER)
card(s, 0.9 * EMU, 3.65 * EMU, 3.45 * EMU, 1.75 * EMU, "_ldap._tcp", "LDAPをTCPで提供", BLUE, body_size=18)
card(s, 4.95 * EMU, 3.65 * EMU, 3.45 * EMU, 1.75 * EMU, "dc._msdcs", "対象ドメインのDCを検索", PURPLE, body_size=17)
card(s, 9.0 * EMU, 3.65 * EMU, 3.45 * EMU, 1.75 * EMU, "結果", "ホスト・ポート\n優先度・重み", GREEN, body_size=18)
bottom_note(s, "Aレコードだけ見て『DNSは正常』と判断しない", ORANGE, 20)

s = add_slide(prs)
header(s, "14", "DOMAIN CONTROLLER", "DC = AD DSを提供するサーバー")
card(s, 0.65 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "保持", "自ドメインの\nディレクトリ情報", BLUE, body_size=22)
card(s, 4.78 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "提供", "認証／LDAP\n時刻同期", TEAL, body_size=22)
card(s, 8.9 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "同期", "他のDCと\n変更を複製", GREEN, body_size=22)
bottom_note(s, "DNS Server役割はDCへの併設が多いが必須ではない", TEXT_MAIN, 21)

s = add_slide(prs)
header(s, "15", "PHYSICAL MODEL", "サイト = ネットワークの『近さ』を表す")
card(s, 0.65 * EMU, 2.0 * EMU, 5.55 * EMU, 3.8 * EMU, "東京サイト", "10.10.0.0/16\nDC-TYO-01／02\n\n高速・安定したLAN", BLUE, body_size=20)
card(s, 7.1 * EMU, 2.0 * EMU, 5.55 * EMU, 3.8 * EMU, "大阪サイト", "10.20.0.0/16\nDC-OSA-01\n\n別拠点のLAN", TEAL, body_size=20)
line(s, 5.75 * EMU, 3.9 * EMU, 7.55 * EMU, 3.9 * EMU, ORANGE, 3.0)
text(s, 5.45 * EMU, 4.1 * EMU, 2.4 * EMU, 0.55 * EMU, [{"text": "WAN／VPN", "size": 18, "bold": True, "color": ORANGE}], align=PP_ALIGN.CENTER)
bottom_note(s, "サイトは組織図ではない。IPサブネットをサイトへ対応付ける", TEXT_MAIN, 20)

s = add_slide(prs)
header(s, "16", "SITE LINK", "サイトリンク = サイト間の接続条件")
card(s, 0.65 * EMU, 2.05 * EMU, 3.5 * EMU, 3.6 * EMU, "経路", "どのサイト同士を\n接続するか", BLUE, body_size=21)
card(s, 4.9 * EMU, 2.05 * EMU, 3.5 * EMU, 3.6 * EMU, "コスト", "複製経路の選択に\n影響する値", ORANGE, body_size=20)
card(s, 9.15 * EMU, 2.05 * EMU, 3.5 * EMU, 3.6 * EMU, "スケジュール", "いつ複製を\n許可するか", GREEN, body_size=21)
bottom_note(s, "近いDCの選択と、遠隔地への複製制御を物理構造で設計", TEXT_MAIN, 20)

s = add_slide(prs)
header(s, "17", "OU OR CONTAINER", "OUは『管理するため』の入れ物")
card(s, 0.7 * EMU, 2.0 * EMU, 5.55 * EMU, 3.8 * EMU, "Organizational Unit（OU）", "オブジェクトを整理\n管理権限を委任\nGPOをリンクできる", GREEN, body_size=21)
card(s, 7.05 * EMU, 2.0 * EMU, 5.55 * EMU, 3.8 * EMU, "Container", "Users／Computersなど\n汎用的な入れ物\nGPOのリンク先ではない", BLUE, body_size=21)
bottom_note(s, "OU設計は会社の部署図コピーではなく『委任と設定』から決める", ORANGE, 20)

s = add_slide(prs)
header(s, "18", "OU ≠ GROUP", "OUは管理場所、グループは権限のまとまり")
card(s, 0.7 * EMU, 2.0 * EMU, 5.55 * EMU, 3.8 * EMU, "OU", "原則として1オブジェクトは\n同時に1つのOU階層へ所属\n\nGPO／委任の対象", GREEN, body_size=19)
card(s, 7.05 * EMU, 2.0 * EMU, 5.55 * EMU, 3.8 * EMU, "セキュリティグループ", "複数グループへ所属できる\n\n共有フォルダー等の\nアクセス権を付与", PURPLE, body_size=19)
bottom_note(s, "『営業部OUだから営業フォルダーへ入れる』ではなく、権限はグループで", TEXT_MAIN, 19)

s = add_slide(prs)
header(s, "19", "GROUP POLICY", "Group Policy Object（GPO）= 設定を一括配布")
steps = [
    ("Local", "端末ローカル", LABEL),
    ("Site", "ネットワーク上のサイト", BLUE),
    ("Domain", "ドメイン全体", PURPLE),
    ("OU", "対象OU", GREEN),
]
for index, (name, body, color) in enumerate(steps):
    x = (0.55 + index * 3.16) * EMU
    card(s, x, 2.05 * EMU, 2.8 * EMU, 3.2 * EMU, name, body, color, body_size=18)
    if index < 3:
        arrow(s, (3.38 + index * 3.16) * EMU, 3.3 * EMU, "→", TEXT_MUTED, 26)
bottom_note(s, "後で処理するGPOが原則優先。強制／継承ブロック／セキュリティフィルターで変わる", ORANGE, 16)

s = add_slide(prs)
header(s, "20", "DOMAIN JOIN", "ドメイン参加 = PCにもアカウントを持たせる")
card(s, 0.65 * EMU, 2.2 * EMU, 3.1 * EMU, 2.75 * EMU, "PC", "ドメインへ\n参加する端末", BLUE, body_size=19)
arrow(s, 3.9 * EMU, 3.25 * EMU, "→", TEXT_MUTED, 32)
card(s, 4.45 * EMU, 2.2 * EMU, 4.0 * EMU, 2.75 * EMU, "ドメイン参加", "必要な権限で\nPCアカウントを作成", ORANGE, body_size=20)
arrow(s, 8.6 * EMU, 3.25 * EMU, "→", TEXT_MUTED, 32)
card(s, 9.15 * EMU, 2.2 * EMU, 3.5 * EMU, 2.75 * EMU, "安全な関係", "端末とドメインを\n結び付ける", GREEN, body_size=19)
bottom_note(s, "前提：DNS・ネットワーク到達性・権限・時刻同期", TEXT_MAIN, 21)

big_statement(
    prs,
    "21",
    "AUTHENTICATION LAYERS",
    "Kerberos／NTLM と",
    "CredSSPは同じ列ではない",
    "Kerberos・NTLMは認証方式。CredSSPは資格情報をリモート先へ委任する仕組み",
    ORANGE,
)

s = add_slide(prs)
header(s, "22", "KERBEROS", "Kerberos = パスワードを毎回渡さず、チケットで進む")
card(s, 0.55 * EMU, 2.2 * EMU, 2.55 * EMU, 2.55 * EMU, "① サインイン", "ユーザーが\n本人確認", BLUE, body_size=19)
arrow(s, 3.15 * EMU, 3.18 * EMU, "→", TEXT_MUTED, 30)
card(s, 3.7 * EMU, 2.2 * EMU, 2.55 * EMU, 2.55 * EMU, "② KDC", "TGTを発行\n＝次の券をもらう券", PURPLE, title_size=20, body_size=16)
arrow(s, 6.3 * EMU, 3.18 * EMU, "→", TEXT_MUTED, 30)
card(s, 6.85 * EMU, 2.2 * EMU, 2.55 * EMU, 2.55 * EMU, "③ サービス券", "TGTを提示し\n対象用の券を取得", ORANGE, body_size=17)
arrow(s, 9.45 * EMU, 3.18 * EMU, "→", TEXT_MUTED, 30)
card(s, 10.0 * EMU, 2.2 * EMU, 2.65 * EMU, 2.55 * EMU, "④ 利用", "ファイル等へ\nアクセス", GREEN, body_size=19)
bottom_note(s, "時刻同期が重要。Kerberosを利用できない条件ではNTLMが残る場合がある", TEXT_MAIN, 17)

s = add_slide(prs)
header(s, "23", "LDAP", "LDAP = ディレクトリへアクセスするプロトコル")
card(s, 0.7 * EMU, 2.0 * EMU, 3.6 * EMU, 3.7 * EMU, "検索", "ユーザーは？\nグループ所属は？", BLUE, body_size=21)
card(s, 4.86 * EMU, 2.0 * EMU, 3.6 * EMU, 3.7 * EMU, "更新", "属性の追加・変更\n権限の範囲内で実行", GREEN, body_size=19)
card(s, 9.02 * EMU, 2.0 * EMU, 3.6 * EMU, 3.7 * EMU, "保護", "署名／TLS／\nChannel Binding", ORANGE, body_size=20)
bottom_note(s, "互換性を監査し、署名・Channel Binding・TLS利用を段階的に強化", TEXT_MAIN, 18)

s = add_slide(prs)
header(s, "24", "DIRECTORY + FILES", "AD DSで複製する主要データは2系統")
card(s, 0.7 * EMU, 2.0 * EMU, 5.55 * EMU, 3.85 * EMU, "NTDS.dit", "ユーザー／PC／グループ\nなどのオブジェクトと属性\n\nAD DSの複製で同期", BLUE, body_size=19)
card(s, 7.05 * EMU, 2.0 * EMU, 5.55 * EMU, 3.85 * EMU, "SYSVOL", "GPOテンプレート\nサインインスクリプト\n\nDFS Replication（DFSR）で同期", GREEN, body_size=18)
bottom_note(s, "ディレクトリデータとSYSVOLは、別の複製系統として確認する", ORANGE, 19)

s = add_slide(prs)
header(s, "25", "REPLICATION", "多くの変更は、どの書き込み可能DCでも受け付ける")
card(s, 0.7 * EMU, 2.2 * EMU, 3.25 * EMU, 2.75 * EMU, "DC-A", "ユーザー属性を\n変更", BLUE, body_size=20)
arrow(s, 4.1 * EMU, 3.25 * EMU, "⇄", TEAL, 34)
card(s, 4.7 * EMU, 2.2 * EMU, 3.25 * EMU, 2.75 * EMU, "DC-B", "グループを\n変更", TEAL, body_size=20)
arrow(s, 8.1 * EMU, 3.25 * EMU, "⇄", TEAL, 34)
card(s, 8.7 * EMU, 2.2 * EMU, 3.95 * EMU, 2.75 * EMU, "DC-C", "変更を受け取り\n複製", GREEN, body_size=20)
bottom_note(s, "マルチマスターが基本。ただし全処理が完全に対等ではない → FSMOへ", ORANGE, 20)

s = add_slide(prs)
header(s, "26", "DELETION LIFECYCLE", "削除状態も、保持するDC間で複製する")
card(s, 0.55 * EMU, 2.15 * EMU, 2.8 * EMU, 2.8 * EMU, "① 削除", "オブジェクトを\n削除", RED, body_size=20)
arrow(s, 3.4 * EMU, 3.25 * EMU, "→", TEXT_MUTED, 30)
card(s, 3.95 * EMU, 2.15 * EMU, 2.8 * EMU, 2.8 * EMU, "② 削除状態", "保持するDC間で\n複製", ORANGE, body_size=19)
arrow(s, 6.8 * EMU, 3.25 * EMU, "→", TEXT_MUTED, 30)
card(s, 7.35 * EMU, 2.15 * EMU, 2.8 * EMU, 2.8 * EMU, "③ 復元可能期間", "Recycle Bin有効時\n多くの属性を保って復元", GREEN, title_size=18, body_size=17)
arrow(s, 10.2 * EMU, 3.25 * EMU, "→", TEXT_MUTED, 30)
card(s, 10.75 * EMU, 2.15 * EMU, 1.9 * EMU, 2.8 * EMU, "④ 消去", "保持後に\n最終消去", PURPLE, body_size=18)
bottom_note(s, "Tombstoneは削除の複製・消去に関係。Recycle Binは復元機能で、バックアップではない", TEXT_MAIN, 16)

s = add_slide(prs)
header(s, "27", "SCHEMA", "スキーマ = ディレクトリの『型定義』")
card(s, 0.7 * EMU, 2.0 * EMU, 3.55 * EMU, 3.75 * EMU, "Class", "User\nComputer\nGroup", BLUE, body_size=21)
card(s, 4.88 * EMU, 2.0 * EMU, 3.55 * EMU, 3.75 * EMU, "Attribute", "displayName\nmail\nmember", TEAL, body_size=21)
card(s, 9.05 * EMU, 2.0 * EMU, 3.55 * EMU, 3.75 * EMU, "Forest-wide", "フォレスト全体で\n1つを共有", ORANGE, body_size=19)
bottom_note(s, "クラスと属性の定義を、フォレスト内のドメインで共有する", TEXT_MAIN, 18)

s = add_slide(prs)
header(s, "28", "FSMO", "5つの操作役割は『担当DCを一つ』に決める")
card(s, 0.6 * EMU, 1.9 * EMU, 5.8 * EMU, 4.2 * EMU, "フォレストに各1つ", "Schema Master\nスキーマ変更を1台で受け付ける\n\nDomain Naming Master\nドメイン追加・削除を調整", PURPLE, body_size=17)
card(s, 6.95 * EMU, 1.9 * EMU, 5.8 * EMU, 4.2 * EMU, "各ドメインに各1つ", "RID Master｜重複しないID部品をDCへ配る\nPDC Emulator｜パスワード変更を優先\nInfrastructure Master｜他ドメイン参照を更新", BLUE, body_size=16)
bottom_note(s, "FSMO = Flexible Single Master Operations。『すべてを処理する主DC』ではない", ORANGE, 18)

s = add_slide(prs)
header(s, "29", "TRUST", "信頼関係 = 別ドメインの認証結果を受け入れる経路")
card(s, 0.75 * EMU, 2.05 * EMU, 4.7 * EMU, 3.5 * EMU, "Domain A", "ユーザーを認証\n『この人です』", BLUE, body_size=21)
arrow(s, 5.7 * EMU, 3.25 * EMU, "→", ORANGE, 38)
card(s, 7.85 * EMU, 2.05 * EMU, 4.7 * EMU, 3.5 * EMU, "Domain B", "認証結果を信頼\n権限を別途確認", TEAL, body_size=21)
bottom_note(s, "信頼があってもアクセス権が自動で付くわけではない", TEXT_MAIN, 22)

s = add_slide(prs)
header(s, "30", "GLOBAL CATALOG", "Global Catalog（GC）= フォレスト横断の索引")
card(s, 0.65 * EMU, 2.2 * EMU, 3.1 * EMU, 2.8 * EMU, "Domain A の情報", "オブジェクトと\n属性", BLUE, body_size=18)
card(s, 5.1 * EMU, 2.2 * EMU, 3.1 * EMU, 2.8 * EMU, "GC機能を持つDC", "自ドメイン：全属性\n他ドメイン：検索用属性", PURPLE, title_size=18, body_size=15)
card(s, 9.55 * EMU, 2.2 * EMU, 3.1 * EMU, 2.8 * EMU, "Domain B の情報", "オブジェクトと\n属性", TEAL, body_size=18)
arrow(s, 3.95 * EMU, 3.25 * EMU, "→", TEXT_MUTED, 30)
arrow(s, 8.4 * EMU, 3.25 * EMU, "←", TEXT_MUTED, 30)
bottom_note(s, "GCは独立サーバーではなくDCへ追加する役割。フォレスト横断検索を助ける", TEXT_MAIN, 18)

big_statement(
    prs,
    "31",
    "RECOVERY",
    "DCを増やしても",
    "複製はバックアップではない",
    "誤削除・侵害・論理破損も複製される。可用性と復旧は別に設計する",
    ORANGE,
)

s = add_slide(prs)
header(s, "32", "HEALTH CHECK", "ADの健全性を確認する代表的な5観点")
checks = [
    ("DNS", "dcdiag /test:dns", TEAL),
    ("複製", "repadmin /replsummary", BLUE),
    ("DC探索", "nltest /dsgetdc:<domain>", PURPLE),
    ("時刻", "w32tm /monitor", ORANGE),
    ("復旧", "System Stateを含む\nAD対応バックアップ\n＋復元テスト", GREEN),
]
for index, (name, command, color) in enumerate(checks):
    if index < 3:
        x = (0.55 + index * 4.15) * EMU
        y = 1.9 * EMU
    else:
        x = (2.65 + (index - 3) * 5.1) * EMU
        y = 4.25 * EMU
    card(s, x, y, 3.75 * EMU, 1.75 * EMU, name, command, color, title_size=20, body_size=15)
bottom_note(s, "運用の目安｜成功表示だけでなく『期待するDC・経路・時刻・復旧性か』を見る", TEXT_MAIN, 17)

s = add_slide(prs)
header(s, "33", "WHAT CHANGED BY 2026", "2026年版で追加して覚える4点")
updates = [
    ("Server 2025", "機能レベルを追加\nその段階のDCは\nServer 2025のみ", BLUE),
    ("SYSVOL複製", "FRS（旧方式）\n→ DFS Replication\n（DFSR）", GREEN),
    ("Kerberos暗号", "旧式RC4の利用元を監査\nAES対応へ移行", ORANGE),
    ("LDAP保護", "署名／TLS／\nChannel Binding\n＝TLS接続との結び付け", PURPLE),
]
for index, (name, body, color) in enumerate(updates):
    x = (0.55 + index * 3.16) * EMU
    card(s, x, 2.0 * EMU, 2.8 * EMU, 3.8 * EMU, name, body, color, title_size=19, body_size=18)
bottom_note(s, "機能レベルはクライアントOSではなく、DCのOSとAD DS機能を規定", TEXT_MAIN, 18)

s = add_slide(prs)
header(s, "34", "SECURITY TRANSITION", "監査と互換性確認から、段階的に強化")
card(s, 0.65 * EMU, 2.05 * EMU, 3.75 * EMU, 3.65 * EMU, "① 監査", "古い暗号RC4の利用元\n保護の弱いLDAP接続元\n\nログで発見", BLUE, body_size=17)
card(s, 4.78 * EMU, 2.05 * EMU, 3.75 * EMU, 3.65 * EMU, "② 互換性を解消", "端末・サービスを更新\nKerberosはAES対応\nLDAPは署名・TLS対応", ORANGE, body_size=17)
card(s, 8.9 * EMU, 2.05 * EMU, 3.75 * EMU, 3.65 * EMU, "③ 段階的な強化", "更新・設定の\n適用状態を確認\n保護を段階的に強化", GREEN, body_size=18)
bottom_note(s, "順番：ログで利用元を特定 → 対応 → 強化。既定値だけに頼らず確認", TEXT_MAIN, 18)

s = add_slide(prs)
header(s, "35", "CONCEPT MAP", "全体はこの線でつながる")
card(s, 0.45 * EMU, 1.85 * EMU, 2.15 * EMU, 1.65 * EMU, "ドメイン", "論理データ区画", PURPLE, title_size=18, body_size=14)
arrow(s, 2.72 * EMU, 2.42 * EMU, "→", TEXT_MUTED, 25)
card(s, 3.25 * EMU, 1.85 * EMU, 2.45 * EMU, 1.65 * EMU, "DNS＋DC Locator", "SRVでDCを発見", TEAL, title_size=16, body_size=14)
arrow(s, 5.82 * EMU, 2.42 * EMU, "→", TEXT_MUTED, 25)
card(s, 6.35 * EMU, 1.85 * EMU, 2.0 * EMU, 1.65 * EMU, "DC", "AD DSを提供", BLUE, title_size=18, body_size=14)
arrow(s, 8.47 * EMU, 2.42 * EMU, "→", TEXT_MUTED, 25)
card(s, 9.0 * EMU, 1.85 * EMU, 3.85 * EMU, 1.65 * EMU, "認証／LDAP", "本人確認・検索・更新", GREEN, title_size=18, body_size=14)
card(s, 0.45 * EMU, 4.0 * EMU, 2.15 * EMU, 1.65 * EMU, "OU", "管理と設定の単位", PURPLE, title_size=18, body_size=14)
arrow(s, 2.72 * EMU, 4.57 * EMU, "→", TEXT_MUTED, 25)
card(s, 3.25 * EMU, 4.0 * EMU, 2.45 * EMU, 1.65 * EMU, "GPO", "設定をリンク", ORANGE, title_size=18, body_size=14)
arrow(s, 5.82 * EMU, 4.57 * EMU, "→", TEXT_MUTED, 25)
card(s, 6.35 * EMU, 4.0 * EMU, 2.0 * EMU, 1.65 * EMU, "ユーザー／PC", "設定を受け取る", GREEN, title_size=16, body_size=14)
card(s, 9.0 * EMU, 4.0 * EMU, 1.45 * EMU, 1.65 * EMU, "Site", "近さ", BLUE, title_size=18, body_size=14)
arrow(s, 10.55 * EMU, 4.57 * EMU, "→", TEXT_MUTED, 25)
card(s, 11.08 * EMU, 4.0 * EMU, 1.77 * EMU, 1.65 * EMU, "近いDC", "選択", TEAL, title_size=16, body_size=14)
bottom_note(s, "DC間の複製：サイトトポロジーとネットワーク接続を考慮して変更を複製", ORANGE, 17)

s = add_slide(prs)
header(s, "36", "NEXT STEP", "次は『実物』を見て、地図と結び付ける")
card(s, 0.65 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "GUI", "ADユーザーとコンピューター\nADサイトとサービス\nDNSマネージャー", BLUE, body_size=18)
card(s, 4.78 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "PowerShell", "Get-ADForest\nGet-ADDomain\nGet-ADDomainController", TEAL, body_size=18)
card(s, 8.9 * EMU, 2.0 * EMU, 3.75 * EMU, 3.65 * EMU, "確認", "SRVレコード\nGPO適用結果\n複製状態", GREEN, body_size=19)
bottom_note(s, "Part 2：構造をGUIとPowerShellで確認する", ORANGE, 23)

s = add_slide(prs)
header(s, "37", "OFFICIAL SOURCES", "今回の正本（基礎）")
sources = [
    ("元動画", "youtube.com/watch?v=lZ8Ps6U_kvY"),
    ("AD DS overview", "learn.microsoft.com/windows-server/identity/ad-ds/.../active-directory-domain-services-overview"),
    ("AD logical model", "learn.microsoft.com/windows-server/identity/ad-ds/plan/understanding-the-active-directory-logical-model"),
    ("AD DS vs Entra ID", "learn.microsoft.com/entra/fundamentals/compare"),
    ("Site topology", "learn.microsoft.com/windows-server/identity/ad-ds/plan/designing-the-site-topology"),
    ("Replication concepts", "learn.microsoft.com/windows-server/identity/ad-ds/.../active-directory-replication-concepts"),
    ("Group Policy processing", "learn.microsoft.com/windows-server/identity/ad-ds/manage/group-policy/group-policy-processing"),
]
for index, (name, url) in enumerate(sources):
    y = (1.55 + index * 0.72) * EMU
    text(s, 0.7 * EMU, y, 3.05 * EMU, 0.5 * EMU, [{"text": name, "size": 15, "bold": True, "color": TEXT_MAIN}])
    text(s, 3.75 * EMU, y, 8.85 * EMU, 0.5 * EMU, [{"text": url, "size": 12, "color": TEAL}])
bottom_note(s, "完全なURLと抽出事実は同梱の source_facts.md に保存", TEXT_MUTED, 17)

s = add_slide(prs)
header(s, "38", "OFFICIAL SOURCES", "今回の正本（2026差分・運用）")
sources = [
    ("Functional levels", "learn.microsoft.com/windows-server/identity/ad-ds/active-directory-functional-levels"),
    ("What's new Server 2025", "learn.microsoft.com/windows-server/get-started/whats-new-windows-server-2025"),
    ("Kerberos RC4", "learn.microsoft.com/windows-server/security/kerberos/detect-remediate-rc4-kerberos"),
    ("CredSSP", "learn.microsoft.com/windows/win32/secauthn/credential-security-support-provider"),
    ("SYSVOL to DFSR", "learn.microsoft.com/windows-server/storage/dfs-replication/migrate-sysvol-to-dfsr"),
    ("FSMO roles", "learn.microsoft.com/windows-server/identity/ad-ds/manage/understand-fsmo-roles"),
    ("DCDiag", "learn.microsoft.com/windows-server/administration/windows-commands/dcdiag"),
]
for index, (name, url) in enumerate(sources):
    y = (1.55 + index * 0.72) * EMU
    text(s, 0.7 * EMU, y, 3.05 * EMU, 0.5 * EMU, [{"text": name, "size": 15, "bold": True, "color": TEXT_MAIN}])
    text(s, 3.75 * EMU, y, 8.85 * EMU, 0.5 * EMU, [{"text": url, "size": 12, "color": TEAL}])
bottom_note(s, "仕様や既定値は更新される。実装時は必ず最新の公式情報を確認", ORANGE, 17)

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(OUT_FILE)
print(OUT_FILE)
